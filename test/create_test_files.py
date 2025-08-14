#!/usr/bin/env python3
"""
Create Test Files for Cloak & Style
Generates comprehensive test files for all supported file types
"""

import os
import csv
import json
from pathlib import Path
from datetime import datetime

def create_csv_test_file():
    """Create a CSV test file with PII data"""
    csv_file = Path("test/test_data.csv")
    
    data = [
        ["Name", "Email", "Phone", "SSN", "Address", "Credit Card"],
        ["John Smith", "john.smith@email.com", "(555) 123-4567", "123-45-6789", "123 Main St, Anytown, ST 12345", "4111-1111-1111-1111"],
        ["Jane Doe", "jane.doe@company.com", "(555) 987-6543", "987-65-4321", "456 Oak Ave, Somewhere, CA 90210", "5555-4444-3333-2222"],
        ["Bob Johnson", "bob.j@business.net", "(555) 456-7890", "456-78-9012", "789 Pine Rd, Elsewhere, TX 54321", "1234-5678-9012-3456"],
        ["Alice Brown", "alice.brown@test.org", "(555) 789-0123", "789-01-2345", "321 Elm St, Nowhere, NY 67890", "9876-5432-1098-7654"],
        ["Charlie Wilson", "charlie.w@demo.com", "(555) 321-6540", "321-65-4321", "654 Maple Dr, Anywhere, FL 13579", "1111-2222-3333-4444"]
    ]
    
    with open(csv_file, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        writer.writerows(data)
    
    print(f"✅ Created CSV test file: {csv_file}")

def create_txt_test_file():
    """Create a text test file with PII data"""
    txt_file = Path("test/test_data.txt")
    
    content = """
CONFIDENTIAL DOCUMENT

Client Information:
- Name: John Smith
- Email: john.smith@company.com
- Phone: (555) 123-4567
- Address: 456 Oak Avenue, Somewhere, CA 90210
- SSN: 987-65-4321
- Credit Card: 5555-4444-3333-2222
- IP Address: 10.0.0.1

Project Details: This project involves sensitive data processing.

Additional Contacts:
- Manager: Jane Doe (jane.doe@email.com)
- Support: support@company.com
- Emergency: (555) 999-8888

Server Information:
- Primary Server: 192.168.1.100
- Backup Server: 172.16.0.50
- Database: db.company.com

Notes: This document contains confidential information that must be protected.
"""
    
    with open(txt_file, 'w', encoding='utf-8') as f:
        f.write(content)
    
    print(f"✅ Created TXT test file: {txt_file}")

def create_docx_test_file():
    """Create a DOCX test file with PII data"""
    try:
        from docx import Document
        from docx.shared import Inches
        
        docx_file = Path("test/test_data.docx")
        
        # Create document
        doc = Document()
        
        # Add title
        title = doc.add_heading('Test Document with PII Data', 0)
        
        # Add paragraph
        doc.add_paragraph('This is a test document with PII data for testing purposes.')
        
        # Add contact information
        doc.add_heading('Contact Information', level=1)
        contact_para = doc.add_paragraph()
        contact_para.add_run('Contact: ').bold = True
        contact_para.add_run('John Smith (john.smith@company.com)')
        
        # Add phone
        phone_para = doc.add_paragraph()
        phone_para.add_run('Phone: ').bold = True
        phone_para.add_run('(555) 123-4567')
        
        # Add SSN
        ssn_para = doc.add_paragraph()
        ssn_para.add_run('SSN: ').bold = True
        ssn_para.add_run('123-45-6789')
        
        # Add table
        doc.add_heading('Employee Data', level=1)
        table = doc.add_table(rows=1, cols=4)
        table.style = 'Table Grid'
        
        # Add header row
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Name'
        hdr_cells[1].text = 'Email'
        hdr_cells[2].text = 'Phone'
        hdr_cells[3].text = 'SSN'
        
        # Add data rows
        row_cells = table.add_row().cells
        row_cells[0].text = 'Jane Doe'
        row_cells[1].text = 'jane.doe@email.com'
        row_cells[2].text = '(555) 987-6543'
        row_cells[3].text = '987-65-4321'
        
        row_cells = table.add_row().cells
        row_cells[0].text = 'Bob Johnson'
        row_cells[1].text = 'bob.j@business.net'
        row_cells[2].text = '(555) 456-7890'
        row_cells[3].text = '456-78-9012'
        
        # Save document
        doc.save(docx_file)
        print(f"✅ Created DOCX test file: {docx_file}")
        
    except ImportError:
        print("⚠️ python-docx not available, skipping DOCX test file")

def create_pptx_test_file():
    """Create a PPTX test file with PII data"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        pptx_file = Path("test/test_data.pptx")
        
        # Create presentation
        prs = Presentation()
        
        # Add slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide1.shapes.title
        subtitle = slide1.placeholders[1]
        
        title.text = "Test Presentation"
        subtitle.text = "John Smith (john.smith@company.com)"
        
        # Add slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide2.shapes.title
        title.text = "Data Summary"
        
        # Add content
        content = slide2.placeholders[1]
        content.text = "SSN: 123-45-6789\nCredit Card: 4111-1111-1111-1111"
        
        # Add speaker notes
        notes_slide = slide2.notes_slide
        notes_slide.notes_text_frame.text = "This slide contains sensitive PII data that needs to be masked."
        
        # Save presentation
        prs.save(pptx_file)
        print(f"✅ Created PPTX test file: {pptx_file}")
        
    except ImportError:
        print("⚠️ python-pptx not available, skipping PPTX test file")

def create_xlsx_test_file():
    """Create an XLSX test file with PII data"""
    try:
        import openpyxl
        from openpyxl.styles import Font
        
        xlsx_file = Path("test/test_data.xlsx")
        
        # Create workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Test Data"
        
        # Add headers
        headers = ["Name", "Email", "Phone", "SSN", "Credit Card", "Address"]
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
        
        # Add data
        data = [
            ["John Smith", "john.smith@company.com", "(555) 123-4567", "123-45-6789", "4111-1111-1111-1111", "123 Main St"],
            ["Jane Doe", "jane.doe@email.com", "(555) 987-6543", "987-65-4321", "5555-4444-3333-2222", "456 Oak Ave"],
            ["Bob Johnson", "bob.j@business.net", "(555) 456-7890", "456-78-9012", "1234-5678-9012-3456", "789 Pine Rd"]
        ]
        
        for row, row_data in enumerate(data, 2):
            for col, value in enumerate(row_data, 1):
                ws.cell(row=row, column=col, value=value)
        
        # Add comment
        ws['A1'].comment = openpyxl.comments.Comment("This cell contains sensitive data", "Test User")
        
        # Add formula
        ws['G1'] = "Formula Test"
        ws['G2'] = "=CONCATENATE(A2,\" \",B2)"
        
        # Save workbook
        wb.save(xlsx_file)
        print(f"✅ Created XLSX test file: {xlsx_file}")
        
    except ImportError:
        print("⚠️ openpyxl not available, skipping XLSX test file")

def create_pdf_test_file():
    """Create a PDF test file with PII data"""
    try:
        import fitz
        
        pdf_file = Path("test/test_data.pdf")
        
        # Create PDF document
        doc = fitz.open()
        
        # Add page 1
        page = doc.new_page()
        page.insert_text((50, 50), "Test PDF Document", fontsize=16)
        page.insert_text((50, 80), "This document contains PII data for testing purposes.", fontsize=12)
        
        # Add PII data
        pii_text = """
Contact Information:
- Name: John Smith
- Email: john.smith@company.com
- Phone: (555) 123-4567
- SSN: 123-45-6789
- Credit Card: 4111-1111-1111-1111
- Address: 456 Oak Avenue, Somewhere, CA 90210
- IP Address: 192.168.1.100

Additional Data:
- Manager: Jane Doe (jane.doe@email.com)
- Support: support@company.com
- Emergency: (555) 999-8888
"""
        
        page.insert_text((50, 120), pii_text, fontsize=10)
        
        # Add page 2
        page2 = doc.new_page()
        page2.insert_text((50, 50), "Page 2 - Additional Information", fontsize=16)
        page2.insert_text((50, 80), "More PII data for testing:", fontsize=12)
        
        additional_text = """
Employee Database:
- Employee ID: 12345
- Department: IT
- Manager: Bob Johnson (bob.j@business.net)
- Phone: (555) 456-7890
- SSN: 456-78-9012
"""
        
        page2.insert_text((50, 120), additional_text, fontsize=10)
        
        # Save PDF
        doc.save(pdf_file)
        doc.close()
        print(f"✅ Created PDF test file: {pdf_file}")
        
    except ImportError:
        print("⚠️ PyMuPDF not available, skipping PDF test file")

def create_md_test_file():
    """Create a Markdown test file with PII data"""
    md_file = Path("test/test_data.md")
    
    content = """# Test Markdown Document

This is a test markdown document containing PII data for testing purposes.

## Contact Information

- **Name:** John Smith
- **Email:** john.smith@company.com
- **Phone:** (555) 123-4567
- **SSN:** 123-45-6789
- **Credit Card:** 4111-1111-1111-1111
- **Address:** 456 Oak Avenue, Somewhere, CA 90210
- **IP Address:** 192.168.1.100

## Employee List

| Name | Email | Phone | SSN |
|------|-------|-------|-----|
| John Smith | john.smith@company.com | (555) 123-4567 | 123-45-6789 |
| Jane Doe | jane.doe@email.com | (555) 987-6543 | 987-65-4321 |
| Bob Johnson | bob.j@business.net | (555) 456-7890 | 456-78-9012 |

## Notes

This document contains sensitive PII data that needs to be masked for privacy compliance.

> **Important:** All PII data must be protected according to company policy.
"""
    
    with open(md_file, 'w', encoding='utf-8') as f:
        f.write(content)
    
    print(f"✅ Created MD test file: {md_file}")

def create_log_test_file():
    """Create a log test file with PII data"""
    log_file = Path("test/test_data.log")
    
    content = """2024-01-15 10:30:15 INFO Application started
2024-01-15 10:30:16 INFO User login: john.smith@company.com
2024-01-15 10:30:17 INFO Processing request from IP: 192.168.1.100
2024-01-15 10:30:18 INFO Database query for SSN: 123-45-6789
2024-01-15 10:30:19 INFO Credit card transaction: 4111-1111-1111-1111
2024-01-15 10:30:20 INFO Contact information: John Smith, (555) 123-4567
2024-01-15 10:30:21 INFO User logout: john.smith@company.com
2024-01-15 10:30:22 INFO Application shutdown
2024-01-15 11:00:00 INFO New session started
2024-01-15 11:00:01 INFO User login: jane.doe@email.com
2024-01-15 11:00:02 INFO Processing request from IP: 172.16.0.50
2024-01-15 11:00:03 INFO Database query for SSN: 987-65-4321
2024-01-15 11:00:04 INFO Contact information: Jane Doe, (555) 987-6543
2024-01-15 11:00:05 INFO User logout: jane.doe@email.com
"""
    
    with open(log_file, 'w', encoding='utf-8') as f:
        f.write(content)
    
    print(f"✅ Created LOG test file: {log_file}")

def create_large_csv_test_file():
    """Create a large CSV test file for streaming testing"""
    large_csv_file = Path("test/large_test_data.csv")
    
    # Create a large CSV file with many rows
    with open(large_csv_file, 'w', newline='', encoding='utf-8') as f:
        writer = csv.writer(f)
        
        # Write header
        writer.writerow(["ID", "Name", "Email", "Phone", "SSN", "Address"])
        
        # Write many rows
        for i in range(1000):  # 1000 rows for testing
            writer.writerow([
                f"ID{i:04d}",
                f"User{i}",
                f"user{i}@test{i}.com",
                f"(555) {i:03d}-{i:04d}",
                f"{i:03d}-{i:02d}-{i:04d}",
                f"{i} Test Street, City{i}, ST {i:05d}"
            ])
    
    print(f"✅ Created large CSV test file: {large_csv_file}")

def main():
    """Create all test files"""
    print("🚀 Creating comprehensive test files for Cloak & Style...")
    print("=" * 60)
    
    # Create test directory if it doesn't exist
    test_dir = Path("test")
    test_dir.mkdir(exist_ok=True)
    
    # Create test files
    create_csv_test_file()
    create_txt_test_file()
    create_docx_test_file()
    create_pptx_test_file()
    create_xlsx_test_file()
    create_pdf_test_file()
    create_md_test_file()
    create_log_test_file()
    create_large_csv_test_file()
    
    print("\n" + "=" * 60)
    print("✅ All test files created successfully!")
    print(f"📁 Test files location: {test_dir.absolute()}")
    
    # List created files
    print("\n📋 Created test files:")
    for file_path in test_dir.glob("*"):
        if file_path.is_file():
            size = file_path.stat().st_size
            print(f"  • {file_path.name} ({size:,} bytes)")

if __name__ == "__main__":
    main()
