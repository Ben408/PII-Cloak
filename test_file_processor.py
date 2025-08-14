#!/usr/bin/env python3
"""
Test script for the File Processor
Tests file processing capabilities with different file types
"""

import sys
import os

# Add the core directory to the path
sys.path.append(os.path.join(os.path.dirname(__file__), 'core'))

from detection_engine import PIIDetectionEngine
from file_processor import FileProcessor

def test_file_processor():
    """Test the file processor with different file types"""
    
    print("🧪 Testing File Processor")
    print("=" * 50)
    
    # Initialize the engine and processor
    engine = PIIDetectionEngine()
    processor = FileProcessor(engine)
    
    print(f"Supported file types: {processor.get_supported_types()}")
    
    # Test 1: Create and process a CSV file
    print("\n📋 Test 1: CSV File Processing")
    print("-" * 30)
    
    test_csv_content = """Name,Email,Phone,Address
John Smith,john.smith@email.com,(555) 123-4567,123 Main St
Jane Doe,jane.doe@company.com,(555) 987-6543,456 Oak Ave
Bob Johnson,bob.j@business.net,(555) 456-7890,789 Pine Rd"""
    
    with open('test_data.csv', 'w', newline='') as f:
        f.write(test_csv_content)
    
    try:
        result = processor.process_file('test_data.csv')
        
        print(f"File: {result.file_info.path}")
        print(f"Type: {result.file_info.file_type}")
        print(f"Size: {result.file_info.size} bytes")
        print(f"Rows: {result.file_info.row_count}")
        print(f"Columns: {result.file_info.column_count}")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        print(f"Entities found: {len(result.entities_found)}")
        
        print("\nDetected entities:")
        for entity in result.entities_found:
            location = getattr(entity, 'location', 'Unknown')
            print(f"  • {entity.entity_type}: '{entity.value}' (confidence: {entity.confidence:.2f}, location: {location})")
        
        print("\nMasked content:")
        print(result.masked_content)
        
    except Exception as e:
        print(f"Error processing CSV: {e}")
    
    finally:
        if os.path.exists('test_data.csv'):
            os.remove('test_data.csv')
    
    # Test 2: Create and process a text file
    print("\n📋 Test 2: Text File Processing")
    print("-" * 30)
    
    test_text_content = """
    CONFIDENTIAL DOCUMENT
    
    Client Information:
    - Name: John Smith
    - Email: john.smith@company.com
    - Phone: (555) 987-6543
    - Address: 456 Oak Avenue, Somewhere, CA 90210
    - SSN: 987-65-4321
    - Credit Card: 5555-4444-3333-2222
    - IP Address: 10.0.0.1
    
    Project Details: This project involves sensitive data processing.
    """
    
    with open('test_data.txt', 'w') as f:
        f.write(test_text_content)
    
    try:
        result = processor.process_file('test_data.txt')
        
        print(f"File: {result.file_info.path}")
        print(f"Type: {result.file_info.file_type}")
        print(f"Size: {result.file_info.size} bytes")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        print(f"Entities found: {len(result.entities_found)}")
        
        print("\nDetected entities:")
        for entity in result.entities_found:
            print(f"  • {entity.entity_type}: '{entity.value}' (confidence: {entity.confidence:.2f})")
        
        print("\nMasked content:")
        print(result.masked_content)
        
    except Exception as e:
        print(f"Error processing text file: {e}")
    
    finally:
        if os.path.exists('test_data.txt'):
            os.remove('test_data.txt')
    
    # Test 3: Create and process a DOCX file
    print("\n📋 Test 3: DOCX File Processing")
    print("-" * 30)
    
    try:
        from docx import Document
        
        doc = Document()
        doc.add_heading('Test Document', 0)
        doc.add_paragraph('This is a test document with PII data.')
        doc.add_paragraph('Contact: John Smith (john.smith@company.com)')
        doc.add_paragraph('Phone: (555) 123-4567')
        doc.add_paragraph('SSN: 123-45-6789')
        
        # Add a table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = 'Name'
        table.cell(0, 1).text = 'Email'
        table.cell(1, 0).text = 'Jane Doe'
        table.cell(1, 1).text = 'jane.doe@email.com'
        
        doc.save('test_data.docx')
        
        result = processor.process_file('test_data.docx')
        
        print(f"File: {result.file_info.path}")
        print(f"Type: {result.file_info.file_type}")
        print(f"Size: {result.file_info.size} bytes")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        print(f"Entities found: {len(result.entities_found)}")
        
        print("\nDetected entities:")
        for entity in result.entities_found:
            print(f"  • {entity.entity_type}: '{entity.value}' (confidence: {entity.confidence:.2f})")
        
        print("\nMasked content:")
        print(result.masked_content)
        
    except Exception as e:
        print(f"Error processing DOCX file: {e}")
    
    finally:
        if os.path.exists('test_data.docx'):
            os.remove('test_data.docx')
    
    # Test 4: Create and process a PPTX file
    print("\n📋 Test 4: PPTX File Processing")
    print("-" * 30)
    
    try:
        from pptx import Presentation
        
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # Title slide
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "Test Presentation"
        
        content = slide.placeholders[1]
        content.text = "Contact: John Smith\nEmail: john.smith@company.com\nPhone: (555) 123-4567"
        
        # Add another slide
        slide_layout = prs.slide_layouts[1]  # Content slide
        slide2 = prs.slides.add_slide(slide_layout)
        
        title2 = slide2.shapes.title
        title2.text = "Data Summary"
        
        content2 = slide2.placeholders[1]
        content2.text = "SSN: 123-45-6789\nCredit Card: 4111-1111-1111-1111"
        
        prs.save('test_data.pptx')
        
        result = processor.process_file('test_data.pptx')
        
        print(f"File: {result.file_info.path}")
        print(f"Type: {result.file_info.file_type}")
        print(f"Size: {result.file_info.size} bytes")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        print(f"Entities found: {len(result.entities_found)}")
        
        print("\nDetected entities:")
        for entity in result.entities_found:
            print(f"  • {entity.entity_type}: '{entity.value}' (confidence: {entity.confidence:.2f})")
        
        print("\nMasked content:")
        print(result.masked_content)
        
    except Exception as e:
        print(f"Error processing PPTX file: {e}")
    
    finally:
        if os.path.exists('test_data.pptx'):
            os.remove('test_data.pptx')
    
    # Test 5: Create and process an XLSX file
    print("\n📋 Test 5: XLSX File Processing")
    print("-" * 30)
    
    try:
        import openpyxl
        
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Test Data"
        
        # Add headers
        ws['A1'] = 'Name'
        ws['B1'] = 'Email'
        ws['C1'] = 'Phone'
        ws['D1'] = 'SSN'
        
        # Add data
        ws['A2'] = 'John Smith'
        ws['B2'] = 'john.smith@company.com'
        ws['C2'] = '(555) 123-4567'
        ws['D2'] = '123-45-6789'
        
        ws['A3'] = 'Jane Doe'
        ws['B3'] = 'jane.doe@email.com'
        ws['C3'] = '(555) 987-6543'
        ws['D3'] = '987-65-4321'
        
        wb.save('test_data.xlsx')
        
        result = processor.process_file('test_data.xlsx')
        
        print(f"File: {result.file_info.path}")
        print(f"Type: {result.file_info.file_type}")
        print(f"Size: {result.file_info.size} bytes")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        print(f"Entities found: {len(result.entities_found)}")
        
        print("\nDetected entities:")
        for entity in result.entities_found:
            location = getattr(entity, 'location', 'Unknown')
            print(f"  • {entity.entity_type}: '{entity.value}' (confidence: {entity.confidence:.2f}, location: {location})")
        
        print("\nMasked content:")
        print(result.masked_content)
        
    except Exception as e:
        print(f"Error processing XLSX file: {e}")
    
    finally:
        if os.path.exists('test_data.xlsx'):
            os.remove('test_data.xlsx')
    
    # Test 6: Create and process a PDF file
    print("\n📋 Test 6: PDF File Processing")
    print("-" * 30)
    
    try:
        import fitz  # PyMuPDF
        
        # Create a simple PDF with test content
        doc = fitz.open()
        page = doc.new_page()
        
        # Add text content with PII
        text_content = """
        CONFIDENTIAL REPORT
        
        Client Information:
        - Name: John Smith
        - Email: john.smith@company.com
        - Phone: (555) 123-4567
        - Address: 123 Main Street, Anytown, ST 12345
        - SSN: 123-45-6789
        - Credit Card: 4111-1111-1111-1111
        - IP Address: 192.168.1.1
        
        Project Details:
        This confidential report contains sensitive information that must be protected.
        All PII should be redacted before sharing with external parties.
        """
        
        page.insert_text((50, 50), text_content)
        doc.save('test_data.pdf')
        doc.close()
        
        result = processor.process_file('test_data.pdf')
        
        print(f"File: {result.file_info.path}")
        print(f"Type: {result.file_info.file_type}")
        print(f"Size: {result.file_info.size} bytes")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        print(f"Entities found: {len(result.entities_found)}")
        
        print("\nDetected entities:")
        for entity in result.entities_found:
            location = getattr(entity, 'location', 'Unknown')
            print(f"  • {entity.entity_type}: '{entity.value}' (confidence: {entity.confidence:.2f}, location: {location})")
        
        print("\nMasked content:")
        print(result.masked_content)
        
    except Exception as e:
        print(f"Error processing PDF file: {e}")
    
    finally:
        if os.path.exists('test_data.pdf'):
            os.remove('test_data.pdf')
    
    print("\n✅ File processor test completed!")

if __name__ == "__main__":
    test_file_processor()
