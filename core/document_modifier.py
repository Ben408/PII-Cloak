#!/usr/bin/env python3
"""
Document Modification Module for Cloak & Style
Handles creating new files with masked content for all supported file types
"""

import os
import csv
import json
from typing import List, Dict, Tuple, Optional, Any
from pathlib import Path
from dataclasses import dataclass

try:
    from .detection_engine import PIIDetectionEngine, PIIEntity, DetectionResult
    from .file_processor import FileProcessor, FileInfo, ProcessingResult
except ImportError:
    from detection_engine import PIIDetectionEngine, PIIEntity, DetectionResult
    from file_processor import FileProcessor, FileInfo, ProcessingResult

@dataclass
class ModificationResult:
    """Result of document modification"""
    original_file: str
    modified_file: str
    file_type: str
    entities_masked: int
    processing_time: float
    errors: List[str]
    html_output: Optional[str] = None
    txt_output: Optional[str] = None
    comments_masked: int = 0
    hyperlinks_processed: int = 0
    tracked_changes_processed: int = 0

class DocumentModifier:
    """Handles full document modification with masked content and advanced features"""
    
    def __init__(self, detection_engine: PIIDetectionEngine):
        self.detection_engine = detection_engine
        self.file_processor = FileProcessor(detection_engine)
        
        # Mapping of file types to modification methods
        self.modification_methods = {
            '.txt': self._modify_text_file,
            '.csv': self._modify_csv_file,
            '.docx': self._modify_docx_file,
            '.pptx': self._modify_pptx_file,
            '.xlsx': self._modify_xlsx_file,
            '.pdf': self._modify_pdf_file,
            '.md': self._modify_text_file,
            '.log': self._modify_text_file
        }
        
        # Advanced modification configuration
        self.config = {
            'generate_html_output': True,
            'generate_txt_output': True,
            'mask_formulas': True,
            'preserve_formatting': True,
            'include_page_breaks': True
        }
    
    def modify_file(self, input_file: str, output_dir: str = None, 
                   mask_format: str = "token") -> ModificationResult:
        """
        Modify a file by creating a new version with masked PII
        
        Args:
            input_file: Path to the input file
            output_dir: Directory to save the modified file (defaults to input directory)
            mask_format: Format for masking ("token" or "asterisk")
            
        Returns:
            ModificationResult with modification details
        """
        import time
        start_time = time.time()
        
        input_path = Path(input_file)
        
        if not input_path.exists():
            raise FileNotFoundError(f"Input file not found: {input_file}")
        
        # Determine output directory
        if output_dir is None:
            output_dir = input_path.parent
        else:
            output_dir = Path(output_dir)
            output_dir.mkdir(parents=True, exist_ok=True)
        
        # Generate output filename
        stem = input_path.stem
        suffix = input_path.suffix.lower()
        output_filename = f"{stem}_masked{suffix}"
        output_path = output_dir / output_filename
        
        # Process the file first
        try:
            processing_result = self.file_processor.process_file(str(input_path))
            
            # Get the modification method
            modifier = self.modification_methods.get(suffix)
            if not modifier:
                raise ValueError(f"Unsupported file type for modification: {suffix}")
            
            # Apply modification
            modification_result = modifier(str(input_path), str(output_path), processing_result, mask_format)
            
            processing_time = time.time() - start_time
            
            return ModificationResult(
                original_file=str(input_path),
                modified_file=str(output_path),
                file_type=suffix,
                entities_masked=len(processing_result.entities_found),
                processing_time=processing_time,
                errors=[],
                html_output=modification_result.get('html_output'),
                txt_output=modification_result.get('txt_output'),
                comments_masked=processing_result.comments_masked,
                hyperlinks_processed=processing_result.hyperlinks_processed,
                tracked_changes_processed=processing_result.tracked_changes_processed
            )
            
        except Exception as e:
            processing_time = time.time() - start_time
            return ModificationResult(
                original_file=str(input_path),
                modified_file="",
                file_type=suffix,
                entities_masked=0,
                processing_time=processing_time,
                errors=[str(e)]
            )
    
    def _modify_text_file(self, input_file: str, output_file: str, 
                         processing_result: ProcessingResult, mask_format: str) -> Dict[str, Any]:
        """Modify text file with masked content"""
        try:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(processing_result.masked_content)
            
            return {}
        except Exception as e:
            raise Exception(f"Error modifying text file: {e}")
    
    def _modify_csv_file(self, input_file: str, output_file: str, 
                        processing_result: ProcessingResult, mask_format: str) -> Dict[str, Any]:
        """Modify CSV file with masked content"""
        try:
            # The masked content is already in CSV format from the processor
            with open(output_file, 'w', encoding='utf-8', newline='') as f:
                f.write(processing_result.masked_content)
            
            return {}
        except Exception as e:
            raise Exception(f"Error modifying CSV file: {e}")
    
    def _modify_docx_file(self, input_file: str, output_file: str, 
                         processing_result: ProcessingResult, mask_format: str) -> Dict[str, Any]:
        """Modify DOCX file with masked content and advanced features"""
        try:
            from docx import Document
            from docx.shared import Inches
            
            # Load the original document
            doc = Document(input_file)
            
            # Create a new document for the masked version
            masked_doc = Document()
            
            # Copy document properties
            masked_doc.core_properties.title = doc.core_properties.title
            masked_doc.core_properties.author = doc.core_properties.author
            masked_doc.core_properties.subject = doc.core_properties.subject
            
            # Process paragraphs with masking
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Get masked text for this paragraph
                    result = self.detection_engine.detect_pii(paragraph.text)
                    masked_text = result.masked_content
                    
                    # Add masked paragraph to new document
                    new_para = masked_doc.add_paragraph(masked_text)
                    
                    # Copy paragraph formatting
                    new_para.style = paragraph.style
                    new_para.alignment = paragraph.alignment
            
            # Process tables
            for table in doc.tables:
                # Create new table with same dimensions
                new_table = masked_doc.add_table(rows=len(table.rows), cols=len(table.columns))
                new_table.style = table.style
                
                for i, row in enumerate(table.rows):
                    for j, cell in enumerate(row.cells):
                        if cell.text.strip():
                            # Get masked text for this cell
                            result = self.detection_engine.detect_pii(cell.text)
                            new_table.cell(i, j).text = result.masked_content
            
            # Process comments if enabled
            if processing_result.comments_masked > 0:
                for comment in doc.comments:
                    # Add comment to new document
                    result = self.detection_engine.detect_pii(comment.text)
                    masked_doc.add_comment(result.masked_content)
            
            # Save the masked document
            masked_doc.save(output_file)
            
            return {}
        except Exception as e:
            raise Exception(f"Error modifying DOCX file: {e}")
    
    def _modify_pptx_file(self, input_file: str, output_file: str, 
                         processing_result: ProcessingResult, mask_format: str) -> Dict[str, Any]:
        """Modify PPTX file with masked content and advanced features"""
        try:
            from pptx import Presentation
            from copy import deepcopy
            
            # Load the original presentation
            prs = Presentation(input_file)
            
            # Create a new presentation
            masked_prs = Presentation()
            
            # Copy slide masters
            for slide_master in prs.slide_masters:
                masked_prs.slide_masters.append(deepcopy(slide_master._element))
            
            # Process slides
            for slide in prs.slides:
                # Create new slide
                new_slide = masked_prs.slides.add_slide(masked_prs.slide_layouts[0])
                
                # Process shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        # Get masked text
                        result = self.detection_engine.detect_pii(shape.text)
                        
                        # Add text box with masked content
                        textbox = new_slide.shapes.add_textbox(
                            shape.left, shape.top, shape.width, shape.height
                        )
                        textbox.text = result.masked_content
                
                # Process speaker notes if enabled
                if processing_result.comments_masked > 0 and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        result = self.detection_engine.detect_pii(notes_text)
                        new_slide.notes_slide.notes_text_frame.text = result.masked_content
            
            # Save the masked presentation
            masked_prs.save(output_file)
            
            return {}
        except Exception as e:
            raise Exception(f"Error modifying PPTX file: {e}")
    
    def _modify_xlsx_file(self, input_file: str, output_file: str, 
                         processing_result: ProcessingResult, mask_format: str) -> Dict[str, Any]:
        """Modify XLSX file with masked content and advanced features"""
        try:
            import openpyxl
            from openpyxl.utils import get_column_letter
            
            # Load the original workbook
            wb = openpyxl.load_workbook(input_file)
            
            # Create a new workbook
            masked_wb = openpyxl.Workbook()
            
            # Process each sheet
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                # Create new sheet
                if sheet_name == masked_wb.active.title:
                    masked_sheet = masked_wb.active
                    masked_sheet.title = sheet_name
                else:
                    masked_sheet = masked_wb.create_sheet(sheet_name)
                
                # Process cells
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value is not None:
                            cell_text = str(cell.value)
                            
                            # Check if cell contains formula
                            if cell_text.startswith('=') and self.config['mask_formulas']:
                                # Mask literals in formulas
                                masked_formula = self._mask_formula_literals(cell_text)
                                masked_sheet[cell.coordinate] = masked_formula
                            else:
                                # Regular text masking
                                result = self.detection_engine.detect_pii(cell_text)
                                masked_sheet[cell.coordinate] = result.masked_content
                
                # Copy formatting
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.has_style:
                            masked_cell = masked_sheet[cell.coordinate]
                            masked_cell.font = cell.font
                            masked_cell.border = cell.border
                            masked_cell.fill = cell.fill
                            masked_cell.number_format = cell.number_format
                            masked_cell.protection = cell.protection
                            masked_cell.alignment = cell.alignment
                
                # Process comments if enabled
                if processing_result.comments_masked > 0:
                    # Process comments in cells
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.comment:
                                result = self.detection_engine.detect_pii(cell.comment.text)
                                masked_sheet[cell.coordinate].comment = openpyxl.comments.Comment(
                                    result.masked_content, cell.comment.author
                                )
            
            # Save the masked workbook
            masked_wb.save(output_file)
            
            return {}
        except Exception as e:
            raise Exception(f"Error modifying XLSX file: {e}")
    
    def _mask_formula_literals(self, formula: str) -> str:
        """Mask literals embedded in Excel formulas while preserving function structure"""
        try:
            # This is a simplified implementation
            # In a full implementation, you would parse the formula and mask only literals
            
            # For now, we'll mask common patterns that might contain PII
            import re
            
            # Mask email addresses in formulas
            formula = re.sub(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', '[EMAIL_XXX]', formula)
            
            # Mask phone numbers in formulas
            formula = re.sub(r'\(\d{3}\)\s*\d{3}-\d{4}', '[PHONE_XXX]', formula)
            
            # Mask SSN in formulas
            formula = re.sub(r'\d{3}-\d{2}-\d{4}', '[SSN_XXX]', formula)
            
            return formula
        except Exception:
            return formula
    
    def _modify_pdf_file(self, input_file: str, output_file: str, 
                        processing_result: ProcessingResult, mask_format: str) -> Dict[str, Any]:
        """Modify PDF file with masked content and generate HTML/TXT output"""
        try:
            import fitz
            
            # Check if it's an image-only PDF
            if processing_result.file_info.is_image_only_pdf:
                raise Exception("Cannot modify image-only PDF - no text content to mask")
            
            # Load the original PDF
            doc = fitz.open(input_file)
            
            # Create a new PDF
            masked_doc = fitz.open()
            
            html_output = ""
            txt_output = ""
            
            # Process each page
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Get masked text for this page
                page_text = page.get_text()
                if page_text.strip():
                    result = self.detection_engine.detect_pii(page_text)
                    masked_text = result.masked_content
                else:
                    masked_text = "[No text content]"
                
                # Create new page
                new_page = masked_doc.new_page(width=page.rect.width, height=page.rect.height)
                
                # Add masked text to new page
                new_page.insert_text((50, 50), masked_text, fontsize=12)
                
                # Add to HTML output
                if self.config['generate_html_output']:
                    html_output += f"<div class='page' id='page-{page_num + 1}'>\n"
                    html_output += f"<h2>Page {page_num + 1}</h2>\n"
                    html_output += f"<div class='content'>{masked_text.replace(chr(10), '<br>')}</div>\n"
                    if self.config['include_page_breaks']:
                        html_output += "<div class='page-break'></div>\n"
                    html_output += "</div>\n"
                
                # Add to TXT output
                if self.config['generate_txt_output']:
                    txt_output += f"Page {page_num + 1}\n"
                    txt_output += "=" * 50 + "\n"
                    txt_output += masked_text + "\n\n"
                    if self.config['include_page_breaks']:
                        txt_output += "\f"  # Form feed for page break
            
            # Save the masked PDF
            masked_doc.save(output_file)
            masked_doc.close()
            doc.close()
            
            # Generate HTML output file
            html_output_file = None
            if self.config['generate_html_output'] and html_output:
                html_output_file = str(Path(output_file).with_suffix('.html'))
                self._generate_html_output(html_output_file, html_output, processing_result)
            
            # Generate TXT output file
            txt_output_file = None
            if self.config['generate_txt_output'] and txt_output:
                txt_output_file = str(Path(output_file).with_suffix('.txt'))
                with open(txt_output_file, 'w', encoding='utf-8') as f:
                    f.write(txt_output)
            
            return {
                'html_output': html_output_file,
                'txt_output': txt_output_file
            }
        except Exception as e:
            raise Exception(f"Error modifying PDF file: {e}")
    
    def _generate_html_output(self, html_file: str, content: str, processing_result: ProcessingResult):
        """Generate HTML output with proper formatting"""
        html_template = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Masked PDF Content</title>
    <style>
        body {{
            font-family: Arial, sans-serif;
            line-height: 1.6;
            margin: 20px;
            background-color: #f5f5f5;
        }}
        .container {{
            max-width: 800px;
            margin: 0 auto;
            background-color: white;
            padding: 20px;
            box-shadow: 0 0 10px rgba(0,0,0,0.1);
        }}
        .page {{
            margin-bottom: 30px;
            padding: 20px;
            border: 1px solid #ddd;
            border-radius: 5px;
        }}
        .page h2 {{
            color: #333;
            border-bottom: 2px solid #007acc;
            padding-bottom: 10px;
        }}
        .content {{
            margin-top: 15px;
            white-space: pre-wrap;
        }}
        .page-break {{
            page-break-after: always;
        }}
        .header {{
            background-color: #007acc;
            color: white;
            padding: 15px;
            margin-bottom: 20px;
            border-radius: 5px;
        }}
        .footer {{
            margin-top: 30px;
            padding: 15px;
            background-color: #f8f9fa;
            border-radius: 5px;
            font-size: 0.9em;
            color: #666;
        }}
        @media print {{
            .page-break {{
                page-break-after: always;
            }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>Masked PDF Content</h1>
            <p>Original file: {processing_result.file_info.path}</p>
            <p>Entities masked: {len(processing_result.entities_found)}</p>
            <p>Processing time: {processing_result.processing_time:.2f} seconds</p>
        </div>
        
        {content}
        
        <div class="footer">
            <p>Generated by Cloak & Style - PII Data Scrubber</p>
            <p>This document contains no raw PII data - only masked tokens for audit purposes.</p>
        </div>
    </div>
</body>
</html>
"""
        
        with open(html_file, 'w', encoding='utf-8') as f:
            f.write(html_template)
    
    def get_supported_modification_types(self) -> List[str]:
        """Get list of file types that support full modification"""
        return list(self.modification_methods.keys())
    
    def can_modify_file(self, file_path: str) -> bool:
        """Check if a file can be fully modified"""
        return Path(file_path).suffix.lower() in self.modification_methods

# Test the document modifier
if __name__ == "__main__":
    # Create a test file
    test_content = """Name,Email,Phone
John Smith,john.smith@email.com,(555) 123-4567
Jane Doe,jane.doe@company.com,(555) 987-6543"""
    
    with open('test_modify.csv', 'w', newline='') as f:
        f.write(test_content)
    
    # Test the modifier
    engine = PIIDetectionEngine()
    modifier = DocumentModifier(engine)
    
    try:
        result = modifier.modify_file('test_modify.csv', mask_format="token")
        
        print("Document modification test:")
        print(f"Original: {result.original_file}")
        print(f"Modified: {result.modified_file}")
        print(f"Type: {result.file_type}")
        print(f"Entities masked: {result.entities_masked}")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        
        if result.errors:
            print(f"Errors: {result.errors}")
        
        # Clean up
        os.remove('test_modify.csv')
        if os.path.exists(result.modified_file):
            os.remove(result.modified_file)
        
    except Exception as e:
        print(f"Error: {e}")
        if os.path.exists('test_modify.csv'):
            os.remove('test_modify.csv')
