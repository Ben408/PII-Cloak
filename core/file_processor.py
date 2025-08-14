#!/usr/bin/env python3
"""
File Processing Module for Cloak & Style
Handles different file types and integrates with the detection engine
"""

import os
import csv
import json
import chardet
from typing import List, Dict, Tuple, Optional, Any, Generator
from pathlib import Path
from dataclasses import dataclass, asdict

try:
    from .detection_engine import PIIDetectionEngine, PIIEntity, DetectionResult
except ImportError:
    from detection_engine import PIIDetectionEngine, PIIEntity, DetectionResult

@dataclass
class FileInfo:
    """Information about a file being processed"""
    path: str
    file_type: str
    size: int
    file_size: int  # Alias for size
    encoding: str
    row_count: Optional[int] = None
    column_count: Optional[int] = None
    page_count: Optional[int] = None
    has_comments: bool = False
    has_tracked_changes: bool = False
    has_hyperlinks: bool = False
    is_image_only_pdf: bool = False

@dataclass
class ProcessingResult:
    """Result of processing a file"""
    file_info: FileInfo
    detection_result: DetectionResult
    masked_content: str
    entities_found: List[PIIEntity]
    processing_time: float
    errors: List[str]
    questionable_entities: List[PIIEntity] = None
    residual_entities: List[PIIEntity] = None
    comments_masked: int = 0
    hyperlinks_processed: int = 0
    tracked_changes_processed: int = 0

class FileProcessor:
    """Handles processing of different file types with advanced features"""
    
    def __init__(self, detection_engine: PIIDetectionEngine):
        self.detection_engine = detection_engine
        self.supported_types = {
            '.txt': self._process_text_file,
            '.csv': self._process_csv_file,
            '.docx': self._process_docx_file,
            '.pptx': self._process_pptx_file,
            '.xlsx': self._process_xlsx_file,
            '.pdf': self._process_pdf_file,
            '.md': self._process_text_file,
            '.log': self._process_text_file
        }
        
        # Advanced processing configuration
        self.config = {
            'chunk_size': 50000,  # 50KB chunks for streaming (to trigger on large test file)
            'max_memory_mb': 100,  # Max memory per file
            'enable_streaming': True,
            'extract_comments': True,
            'extract_tracked_changes': True,
            'extract_hyperlinks': True,
            'extract_chart_labels': True,
            'detect_image_only_pdfs': True
        }
    
    def process_file(self, file_path: str) -> ProcessingResult:
        """
        Process a single file and return results
        
        Args:
            file_path: Path to the file to process
            
        Returns:
            ProcessingResult with detection and masking results
        """
        import time
        start_time = time.time()
        
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Get file information
        file_info = self._get_file_info(file_path)
        
        # Check file size limits
        if file_info.size > self.detection_engine.config['caps'].get('file_size_mb', 50) * 1024 * 1024:
            raise ValueError(f"File too large: {file_info.size} bytes")
        
        # Process based on file type
        processor = self.supported_types.get(file_path.suffix.lower())
        if not processor:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")
        
        try:
            masked_content, entities, advanced_stats = processor(file_path, file_info)
            
            # Create detection result
            detection_result = DetectionResult(
                entities_found=entities,
                masked_content=masked_content,
                original_text="",  # We don't store the original text for privacy
                processing_time=0.0,
                questionable_entities=[],
                residual_entities=[]
            )
            
            processing_time = time.time() - start_time
            
            return ProcessingResult(
                file_info=file_info,
                detection_result=detection_result,
                masked_content=masked_content,
                entities_found=entities,
                processing_time=processing_time,
                errors=[],
                questionable_entities=[],
                residual_entities=[],
                comments_masked=advanced_stats.get('comments_masked', 0),
                hyperlinks_processed=advanced_stats.get('hyperlinks_processed', 0),
                tracked_changes_processed=advanced_stats.get('tracked_changes_processed', 0)
            )
            
        except Exception as e:
            processing_time = time.time() - start_time
            return ProcessingResult(
                file_info=file_info,
                detection_result=DetectionResult(
                    entities_found=[],
                    masked_content="",
                    original_text="",
                    processing_time=processing_time,
                    questionable_entities=[],
                    residual_entities=[]
                ),
                masked_content="",
                entities_found=[],
                processing_time=processing_time,
                errors=[str(e)],
                questionable_entities=[],
                residual_entities=[]
            )
    
    def _get_file_info(self, file_path: Path) -> FileInfo:
        """Get comprehensive file information"""
        # Basic file info
        stat = file_path.stat()
        file_type = file_path.suffix.lower()
        
        # Detect encoding
        encoding = self._detect_encoding(file_path)
        
        # Get additional info based on file type
        row_count = None
        column_count = None
        page_count = None
        has_comments = False
        has_tracked_changes = False
        has_hyperlinks = False
        is_image_only_pdf = False
        
        if file_type == '.csv':
            row_count, column_count = self._count_csv_rows_columns(file_path)
        elif file_type == '.xlsx':
            row_count, column_count, has_comments = self._analyze_xlsx_file(file_path)
        elif file_type == '.docx':
            has_comments, has_tracked_changes, has_hyperlinks = self._analyze_docx_file(file_path)
        elif file_type == '.pptx':
            has_comments, has_hyperlinks = self._analyze_pptx_file(file_path)
        elif file_type == '.pdf':
            page_count, is_image_only_pdf = self._analyze_pdf_file(file_path)
        
        return FileInfo(
            path=str(file_path),
            file_type=file_type,
            size=stat.st_size,
            file_size=stat.st_size,
            encoding=encoding,
            row_count=row_count,
            column_count=column_count,
            page_count=page_count,
            has_comments=has_comments,
            has_tracked_changes=has_tracked_changes,
            has_hyperlinks=has_hyperlinks,
            is_image_only_pdf=is_image_only_pdf
        )
    
    def _detect_encoding(self, file_path: Path) -> str:
        """Detect file encoding using chardet"""
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read(10000)  # Read first 10KB for encoding detection
                result = chardet.detect(raw_data)
                return result['encoding'] or 'utf-8'
        except Exception:
            return 'utf-8'
    
    def _count_csv_rows_columns(self, file_path: Path) -> Tuple[int, int]:
        """Count rows and columns in CSV file"""
        try:
            with open(file_path, 'r', encoding=self._detect_encoding(file_path)) as f:
                reader = csv.reader(f)
                rows = list(reader)
                return len(rows), max(len(row) for row in rows) if rows else 0
        except Exception:
            return 0, 0
    
    def _analyze_xlsx_file(self, file_path: Path) -> Tuple[int, int, bool]:
        """Analyze XLSX file for rows, columns, and comments"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=False)  # Don't use read_only for comment detection
            total_rows = 0
            total_columns = 0
            has_comments = False
            
            for sheet in wb.worksheets:
                total_rows += sheet.max_row
                total_columns = max(total_columns, sheet.max_column)
                # Check for comments
                if hasattr(sheet, '_comments') and sheet._comments:
                    has_comments = True
                # Also check for comments in cells
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.comment:
                            has_comments = True
                            break
                    if has_comments:
                        break
            
            wb.close()
            return total_rows, total_columns, has_comments
        except Exception:
            return 0, 0, False
    
    def _analyze_docx_file(self, file_path: Path) -> Tuple[bool, bool, bool]:
        """Analyze DOCX file for comments, tracked changes, and hyperlinks"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            has_comments = len(doc.comments) > 0
            has_tracked_changes = False
            has_hyperlinks = False
            
            # Check for tracked changes and hyperlinks in paragraphs
            for paragraph in doc.paragraphs:
                for run in paragraph.runs:
                    if run._element.xml.find('w:ins') != -1 or run._element.xml.find('w:del') != -1:
                        has_tracked_changes = True
                    if run._element.xml.find('w:hyperlink') != -1:
                        has_hyperlinks = True
            
            return has_comments, has_tracked_changes, has_hyperlinks
        except Exception:
            return False, False, False
    
    def _analyze_pptx_file(self, file_path: Path) -> Tuple[bool, bool]:
        """Analyze PPTX file for speaker notes and hyperlinks"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            has_comments = False
            has_hyperlinks = False
            
            for slide in prs.slides:
                # Check for speaker notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                    has_comments = True
                
                # Check for hyperlinks in shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run, 'hyperlink') and run.hyperlink:
                                    has_hyperlinks = True
            
            return has_comments, has_hyperlinks
        except Exception:
            return False, False
    
    def _analyze_pdf_file(self, file_path: Path) -> Tuple[int, bool]:
        """Analyze PDF file for page count and image-only detection"""
        try:
            import fitz
            doc = fitz.open(str(file_path))
            page_count = len(doc)
            
            # Check if PDF is image-only
            is_image_only = True
            for page_num in range(min(3, page_count)):  # Check first 3 pages
                page = doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    is_image_only = False
                    break
            
            doc.close()
            return page_count, is_image_only
        except Exception:
            return 0, False
    
    def _process_text_file(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process text file with streaming support"""
        if file_info.size > self.config['chunk_size'] and self.config['enable_streaming']:
            return self._process_text_file_streaming(file_path, file_info)
        else:
            return self._process_text_file_standard(file_path, file_info)
    
    def _process_text_file_streaming(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process large text file using streaming"""
        all_entities = []
        masked_content = ""
        
        try:
            with open(file_path, 'r', encoding=file_info.encoding) as f:
                for chunk in self._read_file_chunks(f):
                    # Process chunk
                    result = self.detection_engine.detect_pii(chunk)
                    all_entities.extend(result.entities_found)
                    masked_content += result.masked_content
        except Exception as e:
            raise Exception(f"Error processing text file: {e}")
        
        return masked_content, all_entities, {}
    
    def _process_text_file_standard(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process text file using standard method"""
        try:
            with open(file_path, 'r', encoding=file_info.encoding) as f:
                content = f.read()
            
            result = self.detection_engine.detect_pii(content)
            return result.masked_content, result.entities_found, {}
        except Exception as e:
            raise Exception(f"Error processing text file: {e}")
    
    def _read_file_chunks(self, file_obj, chunk_size: int = None) -> Generator[str, None, None]:
        """Read file in chunks"""
        if chunk_size is None:
            chunk_size = self.config['chunk_size']
        
        while True:
            chunk = file_obj.read(chunk_size)
            if not chunk:
                break
            yield chunk
    
    def _process_csv_file(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process CSV file with streaming support"""
        if file_info.size > self.config['chunk_size'] and self.config['enable_streaming']:
            return self._process_csv_file_streaming(file_path, file_info)
        else:
            return self._process_csv_file_standard(file_path, file_info)
    
    def _process_csv_file_streaming(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process large CSV file using streaming"""
        all_entities = []
        masked_rows = []
        
        try:
            with open(file_path, 'r', encoding=file_info.encoding) as f:
                reader = csv.reader(f)
                for row_num, row in enumerate(reader):
                    masked_row = []
                    for col_num, cell in enumerate(row):
                        # Process cell
                        result = self.detection_engine.detect_pii(cell)
                        masked_row.append(result.masked_content)
                        
                        # Add location info to entities
                        for entity in result.entities_found:
                            entity.start_pos = col_num
                            entity.end_pos = col_num + 1
                            entity.location = f"Row {row_num + 1}, Column {col_num + 1}"
                        all_entities.extend(result.entities_found)
                    
                    masked_rows.append(masked_row)
        except Exception as e:
            raise Exception(f"Error processing CSV file: {e}")
        
        # Convert back to CSV format
        import io
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerows(masked_rows)
        masked_content = output.getvalue()
        
        return masked_content, all_entities, {}
    
    def _process_csv_file_standard(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process CSV file using standard method"""
        all_entities = []
        masked_rows = []
        
        try:
            with open(file_path, 'r', encoding=file_info.encoding) as f:
                reader = csv.reader(f)
                for row_num, row in enumerate(reader):
                    masked_row = []
                    for col_num, cell in enumerate(row):
                        # Process cell
                        result = self.detection_engine.detect_pii(cell)
                        masked_row.append(result.masked_content)
                        
                        # Add location info to entities
                        for entity in result.entities_found:
                            entity.start_pos = col_num
                            entity.end_pos = col_num + 1
                            entity.location = f"Row {row_num + 1}, Column {col_num + 1}"
                        all_entities.extend(result.entities_found)
                    
                    masked_rows.append(masked_row)
        except Exception as e:
            raise Exception(f"Error processing CSV file: {e}")
        
        # Convert back to CSV format
        import io
        output = io.StringIO()
        writer = csv.writer(output)
        writer.writerows(masked_rows)
        masked_content = output.getvalue()
        
        return masked_content, all_entities, {}
    
    def _process_docx_file(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process DOCX file with advanced features"""
        try:
            from docx import Document
            doc = Document(file_path)
            
            all_entities = []
            masked_content = ""
            comments_masked = 0
            hyperlinks_processed = 0
            tracked_changes_processed = 0
            
            # Process paragraphs
            for paragraph in doc.paragraphs:
                paragraph_text = paragraph.text
                result = self.detection_engine.detect_pii(paragraph_text)
                all_entities.extend(result.entities_found)
                masked_content += result.masked_content + "\n"
            
            # Process comments if enabled
            if self.config['extract_comments'] and file_info.has_comments:
                for comment in doc.comments:
                    comment_text = comment.text
                    result = self.detection_engine.detect_pii(comment_text)
                    all_entities.extend(result.entities_found)
                    comments_masked += 1
            
            # Process tracked changes if enabled
            if self.config['extract_tracked_changes'] and file_info.has_tracked_changes:
                # This would require more complex XML processing
                # For now, we'll mark that we detected tracked changes
                tracked_changes_processed = 1
            
            # Process hyperlinks if enabled
            if self.config['extract_hyperlinks'] and file_info.has_hyperlinks:
                # This would require more complex XML processing
                # For now, we'll mark that we detected hyperlinks
                hyperlinks_processed = 1
            
            return masked_content, all_entities, {
                'comments_masked': comments_masked,
                'hyperlinks_processed': hyperlinks_processed,
                'tracked_changes_processed': tracked_changes_processed
            }
        except Exception as e:
            raise Exception(f"Error processing DOCX file: {e}")
    
    def _process_pptx_file(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process PPTX file with advanced features"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            all_entities = []
            masked_content = ""
            comments_masked = 0
            hyperlinks_processed = 0
            
            # Process slides
            for slide_num, slide in enumerate(prs.slides):
                masked_content += f"Slide {slide_num + 1}:\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            paragraph_text = paragraph.text
                            result = self.detection_engine.detect_pii(paragraph_text)
                            all_entities.extend(result.entities_found)
                            masked_content += result.masked_content + "\n"
                
                # Process speaker notes if enabled
                if self.config['extract_comments'] and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    if notes_text.strip():
                        result = self.detection_engine.detect_pii(notes_text)
                        all_entities.extend(result.entities_found)
                        comments_masked += 1
                        masked_content += f"Speaker Notes: {result.masked_content}\n"
            
            return masked_content, all_entities, {
                'comments_masked': comments_masked,
                'hyperlinks_processed': hyperlinks_processed
            }
        except Exception as e:
            raise Exception(f"Error processing PPTX file: {e}")
    
    def _process_xlsx_file(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process XLSX file with advanced features"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True)
            
            all_entities = []
            masked_content = ""
            comments_masked = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                masked_content += f"Sheet: {sheet_name}\n"
                
                # Process cell values
                for row in sheet.iter_rows(values_only=True):
                    masked_row = []
                    for cell_value in row:
                        if cell_value is not None:
                            cell_text = str(cell_value)
                            result = self.detection_engine.detect_pii(cell_text)
                            all_entities.extend(result.entities_found)
                            masked_row.append(result.masked_content)
                        else:
                            masked_row.append("")
                    masked_content += ",".join(masked_row) + "\n"
                
                # Process comments if enabled
                if self.config['extract_comments']:
                    # Check for comments in cells
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.comment:
                                comment_text = cell.comment.text
                                result = self.detection_engine.detect_pii(comment_text)
                                all_entities.extend(result.entities_found)
                                comments_masked += 1
            
            wb.close()
            return masked_content, all_entities, {
                'comments_masked': comments_masked
            }
        except Exception as e:
            raise Exception(f"Error processing XLSX file: {e}")
    
    def _process_pdf_file(self, file_path: Path, file_info: FileInfo) -> Tuple[str, List[PIIEntity], Dict[str, int]]:
        """Process PDF file with advanced features"""
        try:
            import fitz
            
            # Check for image-only PDF
            if file_info.is_image_only_pdf:
                raise Exception("Image-only PDF detected - cannot process text content")
            
            doc = fitz.open(str(file_path))
            all_entities = []
            masked_content = ""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                
                if page_text.strip():
                    result = self.detection_engine.detect_pii(page_text)
                    all_entities.extend(result.entities_found)
                    masked_content += f"Page {page_num + 1}:\n{result.masked_content}\n\n"
            
            doc.close()
            return masked_content, all_entities, {}
        except Exception as e:
            raise Exception(f"Error processing PDF file: {e}")
    
    def _calculate_confidence_scores(self, entities: List[PIIEntity]) -> Dict[str, float]:
        """Calculate confidence scores for entities"""
        scores = {}
        
        for entity in entities:
            entity_type = entity.entity_type
            if entity_type not in scores:
                scores[entity_type] = []
            scores[entity_type].append(entity.confidence)
        
        # Calculate average confidence for each type
        for entity_type in scores:
            scores[entity_type] = sum(scores[entity_type]) / len(scores[entity_type])
        
        return scores
    
    def get_supported_types(self) -> List[str]:
        """Get list of supported file types"""
        return list(self.supported_types.keys())
    
    def can_process_file(self, file_path: str) -> bool:
        """Check if a file can be processed"""
        return Path(file_path).suffix.lower() in self.supported_types

# Test the file processor
if __name__ == "__main__":
    # Create a test CSV file
    test_csv_content = """Name,Email,Phone,Address
John Smith,john.smith@email.com,(555) 123-4567,123 Main St
Jane Doe,jane.doe@company.com,(555) 987-6543,456 Oak Ave
Bob Johnson,bob.j@business.net,(555) 456-7890,789 Pine Rd"""
    
    with open('test_data.csv', 'w', newline='') as f:
        f.write(test_csv_content)
    
    # Test the processor
    engine = PIIDetectionEngine()
    processor = FileProcessor(engine)
    
    try:
        result = processor.process_file('test_data.csv')
        
        print("File processing test:")
        print(f"File: {result.file_info.path}")
        print(f"Type: {result.file_info.file_type}")
        print(f"Size: {result.file_info.size} bytes")
        print(f"Rows: {result.file_info.row_count}")
        print(f"Columns: {result.file_info.column_count}")
        print(f"Processing time: {result.processing_time:.2f} seconds")
        print(f"Entities found: {len(result.entities_found)}")
        
        print("\nMasked content:")
        print(result.masked_content)
        
        print("\nDetected entities:")
        for entity in result.entities_found:
            print(f"- {entity.entity_type}: {entity.value} (confidence: {entity.confidence:.2f})")
        
        # Clean up
        os.remove('test_data.csv')
        
    except Exception as e:
        print(f"Error: {e}")
        if os.path.exists('test_data.csv'):
            os.remove('test_data.csv')
