#!/usr/bin/env python3
"""
Cloak & Style - PII Data Scrubber UI (Fixed Version)
Professional desktop application for PII masking
Fixed version resolving layout overlaps and functionality issues
"""

import sys
import os
from pathlib import Path
from datetime import datetime
from PySide6.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, 
                               QHBoxLayout, QTabWidget, QLabel, QPushButton, 
                               QTextEdit, QProgressBar, QFileDialog, QMessageBox,
                               QTableWidget, QTableWidgetItem, QHeaderView, QCheckBox,
                               QComboBox, QGroupBox, QSplitter, QFrame, QScrollArea,
                               QLineEdit, QListWidget, QListWidgetItem, QMenuBar,
                               QStatusBar, QToolBar, QSpacerItem, QSizePolicy)
from PySide6.QtCore import Qt, Signal, QThread, QTimer, QSize
from PySide6.QtGui import QAction, QDragEnterEvent, QDropEvent, QFont, QPalette, QColor, QIcon

# Add the core directory to the path
sys.path.append(os.path.join(os.path.dirname(__file__), 'core'))

try:
    from detection_engine import PIIDetectionEngine, PIIEntity, DetectionResult
    from document_modifier import DocumentModifier, ModificationResult
    from report_generator import ReportGenerator
    from performance_optimizer import LaptopOptimizedProcessor, PerformanceCaps
except ImportError as e:
    print(f"Import error: {e}")
    # Create dummy classes for testing
    class PIIDetectionEngine:
        def __init__(self): pass
        def detect_pii(self, text): return None
    class DocumentModifier:
        def __init__(self, engine): pass
        def modify_file(self, file, output_dir): return None
    class ReportGenerator:
        def __init__(self): pass
        def generate_html_report(self, results, config, path): return path
        def generate_json_report(self, results, config, path): return path
        def generate_csv_findings(self, results, path): return path
    class LaptopOptimizedProcessor:
        def __init__(self): pass
        def process_files(self, files, func): return []

brand_keywords = []


class DarkTheme:
    """Dark theme colors matching the sophisticated screenshots"""
    BACKGROUND = "#1e1e1e"
    SURFACE = "#2d2d30"
    SURFACE_LIGHT = "#3e3e42"
    TEXT_PRIMARY = "#ffffff"
    TEXT_SECONDARY = "#cccccc"
    ACCENT = "#0078d4"
    SUCCESS = "#107c10"
    WARNING = "#ff8c00"
    ERROR = "#d13438"
    BORDER = "#555555"
    HIGHLIGHT = "#264f78"

class FixedProcessingOptions(QWidget):
    """Fixed processing options with proper layout spacing"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setup_ui()
    
    def setup_ui(self):
        layout = QVBoxLayout()
        layout.setSpacing(15)  # Reduced spacing
        layout.setContentsMargins(0, 0, 0, 0)  # No margins
        
        # Review Queue Section
        review_group = QGroupBox("Review Queue")
        review_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 12px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                margin-top: 10px;
                padding-top: 10px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 5px 0 5px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        review_layout = QVBoxLayout()
        review_layout.setSpacing(8)  # Reduced spacing
        
        self.enable_review = QCheckBox("Enable Review Queue")
        self.enable_review.setStyleSheet(f"""
            QCheckBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-size: 11px;
                spacing: 6px;
            }}
            QCheckBox::indicator {{
                width: 16px;
                height: 16px;
                border: 2px solid {DarkTheme.BORDER};
                border-radius: 3px;
                background-color: {DarkTheme.SURFACE_LIGHT};
            }}
            QCheckBox::indicator:checked {{
                background-color: {DarkTheme.ACCENT};
                border-color: {DarkTheme.ACCENT};
            }}
        """)
        review_layout.addWidget(self.enable_review)
        
        self.dry_run = QCheckBox("Dry Run (Preview Only)", tristate=False)
        self.dry_run.setStyleSheet(f"""
            QCheckBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-size: 11px;
                spacing: 6px;
            }}
            QCheckBox::indicator {{
                width: 16px;
                height: 16px;
                border: 2px solid {DarkTheme.BORDER};
                border-radius: 3px;
                background-color: {DarkTheme.SURFACE_LIGHT};
            }}
            QCheckBox::indicator:checked {{
                background-color: {DarkTheme.ACCENT};
                border-color: {DarkTheme.ACCENT};
            }}
        """)
        review_layout.addWidget(self.dry_run)
        
        self.process_subfolders = QCheckBox("Process Subfolders")
        self.process_subfolders.setStyleSheet(f"""
            QCheckBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-size: 11px;
                spacing: 6px;
            }}
            QCheckBox::indicator {{
                width: 16px;
                height: 16px;
                border: 2px solid {DarkTheme.BORDER};
                border-radius: 3px;
                background-color: {DarkTheme.SURFACE_LIGHT};
            }}
            QCheckBox::indicator:checked {{
                background-color: {DarkTheme.ACCENT};
                border-color: {DarkTheme.ACCENT};
            }}
        """)
        review_layout.addWidget(self.process_subfolders)
        
        review_group.setLayout(review_layout)
        layout.addWidget(review_group)
        
        # Masking Options Section
        masking_group = QGroupBox("Masking Options")
        masking_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 12px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                margin-top: 10px;
                padding-top: 10px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 5px 0 5px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        masking_layout = QVBoxLayout()
        masking_layout.setSpacing(8)
        
        mask_label = QLabel("Mask Format:")
        mask_label.setStyleSheet(f"color: {DarkTheme.TEXT_PRIMARY}; font-size: 11px; font-weight: 500;")
        masking_layout.addWidget(mask_label)
        
        self.mask_format = QComboBox()
        self.mask_format.addItems(["Token Format [TYPE_###]", "Asterisk Format [***]"])
        self.mask_format.setStyleSheet(f"""
            QComboBox {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 4px;
                padding: 6px 10px;
                color: {DarkTheme.TEXT_PRIMARY};
                font-size: 11px;
                min-height: 18px;
            }}
            QComboBox::drop-down {{
                border: none;
                width: 16px;
            }}
            QComboBox::down-arrow {{
                image: none;
                border-left: 4px solid transparent;
                border-right: 4px solid transparent;
                border-top: 4px solid {DarkTheme.TEXT_PRIMARY};
                margin-right: 6px;
            }}
        """)
        masking_layout.addWidget(self.mask_format)
        
        masking_group.setLayout(masking_layout)
        layout.addWidget(masking_group)
        
        # Brand Keywords Section
        brand_group = QGroupBox("Brand Keywords")
        brand_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 12px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                margin-top: 10px;
                padding-top: 10px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 5px 0 5px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        brand_layout = QVBoxLayout()
        brand_layout.setSpacing(8)
        
        button_layout = QHBoxLayout()
        button_layout.setSpacing(8)
        
        upload_btn = QPushButton("Upload Keyword List")
        upload_btn.clicked.connect(self.keywords_manual)
        upload_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {DarkTheme.ACCENT};
                color: white;
                border: none;
                border-radius: 4px;
                padding: 8px 12px;
                font-weight: 600;
                font-size: 11px;
                min-height: 18px;
            }}
            QPushButton:hover {{
                background-color: #106ebe;
            }}
        """)
        button_layout.addWidget(upload_btn)
        
        clear_btn = QPushButton("Clear")
        clear_btn.clicked.connect(self.clear_keywords)
        clear_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                color: {DarkTheme.TEXT_PRIMARY};
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 4px;
                padding: 8px 12px;
                font-weight: 600;
                font-size: 11px;
                min-height: 18px;
            }}
            QPushButton:hover {{
                background-color: {DarkTheme.BORDER};
            }}
        """)
        button_layout.addWidget(clear_btn)
        file_upload_btn = QPushButton("Upload Keyword File")
        file_upload_btn.clicked.connect(self.keyword_list)
        file_upload_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                color: {DarkTheme.TEXT_PRIMARY};
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 4px;
                padding: 8px 12px;
                font-weight: 600;
                font-size: 11px;
                min-height: 18px;
            }}
            QPushButton:hover {{
                background-color: {DarkTheme.BORDER};
            }}
        """)
        button_layout.addWidget(file_upload_btn)
        brand_layout.addLayout(button_layout)
        
        self.keyword_text = QTextEdit()
        self.keyword_text.setMaximumHeight(80)  # Reduced height
        self.keyword_text.setPlaceholderText("Enter keywords separated by a comma or upload a file...")
        self.keyword_text.setStyleSheet(f"""
            QTextEdit {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 4px;
                color: {DarkTheme.TEXT_PRIMARY};
                padding: 8px;
                font-size: 11px;
                line-height: 1.4;
            }}
        """)
        brand_layout.addWidget(self.keyword_text)
        
        brand_group.setLayout(brand_layout)
        layout.addWidget(brand_group)
        
        # Performance Limits Section
        limits_group = QGroupBox("Performance Limits")
        limits_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 12px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                margin-top: 10px;
                padding-top: 10px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 10px;
                padding: 0 5px 0 5px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        limits_layout = QVBoxLayout()
        limits_layout.setSpacing(4)
        
        limits_text = [
            "PDF Pages: ≤ 100 pages",
            "PDF Size: ≤ 10 MB", 
            "Excel/CSV Rows: ≤ 100,000 rows"
        ]
        
        for text in limits_text:
            label = QLabel(text)
            label.setStyleSheet(f"color: {DarkTheme.TEXT_SECONDARY}; font-size: 10px; padding: 1px 0;")
            limits_layout.addWidget(label)
        
        limits_group.setLayout(limits_layout)
        layout.addWidget(limits_group)
        
        layout.addStretch()
        self.setLayout(layout)

    def keyword_list(self):
        files, _ = QFileDialog.getOpenFileNames(
                self,
                "Select Files",
                "",
                "*.txt"
            )
        if files:
            self.add_keywords(files, "file")

    def keywords_manual(self):
        terms = self.keyword_text.toPlainText().split(',')
        self.add_keywords(terms, "manual")

    def clear_keywords(self):
        global brand_keywords
        brand_keywords = []
        print("Brand keywords cleared. Current keywords:", brand_keywords)

    def add_keywords(self, files, type):
        if type == "file":
            if len(files) > 1:
                for file in files:
                    with open(file, 'r') as file:
                        file = file.split(',').strip()
                        for word in file:
                            if word not in brand_keywords:
                                brand_keywords.append(word)
        else:
            for word in files:
                if word.strip() not in brand_keywords:
                    brand_keywords.append(word.strip())
        print("Current Brand Keywords: ", brand_keywords)


    
class FixedFilePreview(QWidget):
    """Fixed file preview with working functionality"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setup_ui()
    
    def setup_ui(self):
        layout = QVBoxLayout()
        layout.setSpacing(10)
        
        # Tab widget
        self.tab_widget = QTabWidget()
        self.tab_widget.setStyleSheet(f"""
            QTabWidget::pane {{
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                background-color: {DarkTheme.SURFACE};
                margin-top: -1px;
            }}
            QTabBar::tab {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                color: {DarkTheme.TEXT_PRIMARY};
                padding: 8px 16px;
                border: 1px solid {DarkTheme.BORDER};
                border-bottom: none;
                border-top-left-radius: 6px;
                border-top-right-radius: 6px;
                font-weight: 500;
                font-size: 11px;
                min-width: 80px;
            }}
            QTabBar::tab:selected {{
                background-color: {DarkTheme.SURFACE};
                border-bottom: 1px solid {DarkTheme.SURFACE};
                color: {DarkTheme.ACCENT};
            }}
            QTabBar::tab:hover:!selected {{
                background-color: {DarkTheme.BORDER};
            }}
        """)
        
        # Text Preview Tab
        self.text_preview = QTextEdit()
        self.text_preview.setReadOnly(True)
        self.text_preview.setPlaceholderText("No file selected")
        self.text_preview.setStyleSheet(f"""
            QTextEdit {{
                background-color: {DarkTheme.SURFACE};
                color: {DarkTheme.TEXT_PRIMARY};
                border: none;
                padding: 10px;
                font-family: 'Consolas', 'Monaco', monospace;
                font-size: 11px;
                line-height: 1.4;
            }}
        """)
        self.tab_widget.addTab(self.text_preview, "Text Preview")
        
        # Table Preview Tab
        self.table_preview = QTextEdit()
        self.table_preview.setReadOnly(True)
        self.table_preview.setPlaceholderText("No table data available")
        self.table_preview.setStyleSheet(f"""
            QTextEdit {{
                background-color: {DarkTheme.SURFACE};
                color: {DarkTheme.TEXT_PRIMARY};
                border: none;
                padding: 10px;
                font-family: 'Consolas', 'Monaco', monospace;
                font-size: 11px;
                line-height: 1.4;
            }}
        """)
        self.tab_widget.addTab(self.table_preview, "Table Preview")
        
        # Detected Entities Tab
        self.entities_preview = QTextEdit()
        self.entities_preview.setReadOnly(True)
        self.entities_preview.setPlaceholderText("No entities detected")
        self.entities_preview.setStyleSheet(f"""
            QTextEdit {{
                background-color: {DarkTheme.SURFACE};
                color: {DarkTheme.TEXT_PRIMARY};
                border: none;
                padding: 10px;
                font-family: 'Consolas', 'Monaco', monospace;
                font-size: 11px;
                line-height: 1.4;
            }}
        """)
        self.tab_widget.addTab(self.entities_preview, "Detected Entities")
        
        layout.addWidget(self.tab_widget)
        self.setLayout(layout)
    
    def load_file_content(self, file_path):
        """Load and display file content in preview"""
        try:
            if not os.path.exists(file_path):
                self.text_preview.setText("File not found")
                return
            
            content = ""
            table_content = ""
            # Handle different file types
            if file_path.endswith((".txt", ".md", ".csv")):
                # Read text-based files
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                table_content = content[:10000]  # Limit to first 10k chars
                
            elif file_path.endswith(".xlsx"):
                # Read XLSX files
                try:
                    from openpyxl import load_workbook
                    wb = load_workbook(filename=file_path)
                    
                    # Get the first worksheet
                    ws = wb.active
                    
                    # Convert to text format for display
                    rows_data = []
                    for row in ws.iter_rows(values_only=True):
                        row_data = []
                        for cell in row:
                            if cell is None:
                                row_data.append("")
                            else:
                                row_data.append(str(cell))
                        rows_data.append("\t".join(row_data))
                    
                    content = "\n".join(rows_data)
                    table_content = content[:10000]  # Limit to first 10k chars
                    
                except ImportError:
                    # Fallback if openpyxl is not available
                    content = "XLSX file detected but openpyxl library not available for preview"
                    table_content = content
                except Exception as e:
                    content = f"Error reading XLSX file: {str(e)}"
                    table_content = content
                    
            else:
                content = "File format not supported for preview"
                table_content = content
            
            # Display in text preview
            if file_path.endswith((".txt",".md", ".csv")):
                self.text_preview.setText(content[:10000])  # Limit to first 10k chars
            elif file_path.endswith(".xlsx"):
                self.text_preview.setText("File displayed on Table Preview")
            else:
                self.text_preview.setText("File format not supported for preview")
            
            # Display in table preview
            if file_path.endswith(".xlsx"):
                self.table_preview.setText(table_content)
            elif file_path.endswith((".txt",".md", ".csv")):
                self.table_preview.setText("File displayed on Text Preview")
            else:
                self.table_preview.setText("File format not supported for preview")
            
            # Show file info in entities tab
            file_info = f"File: {os.path.basename(file_path)}\n"
            file_info += f"Size: {len(content)} characters\n"
            file_info += f"Path: {file_path}\n\n"
            file_info += "No PII entities detected yet.\n"
            file_info += "Run processing to detect entities."
            self.entities_preview.setText(file_info)
            
        except Exception as e:
            error_msg = f"Error loading file:\n{str(e)}"
            self.text_preview.setText(error_msg)
            self.table_preview.setText(error_msg)
            self.entities_preview.setText(error_msg)

class FixedProcessingStatus(QWidget):
    """Fixed processing status with proper layout"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setup_ui()
    
    def setup_ui(self):
        layout = QVBoxLayout()
        layout.setSpacing(10)
        
        # Overall Progress Section
        overall_label = QLabel("Overall Progress:")
        overall_label.setStyleSheet(f"color: {DarkTheme.TEXT_PRIMARY}; font-weight: 600; font-size: 12px;")
        layout.addWidget(overall_label)
        
        self.overall_progress = QProgressBar()
        self.overall_progress.setStyleSheet(f"""
            QProgressBar {{
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                text-align: center;
                background-color: {DarkTheme.SURFACE_LIGHT};
                height: 20px;
                font-weight: 500;
                font-size: 11px;
            }}
            QProgressBar::chunk {{
                background-color: {DarkTheme.ACCENT};
                border-radius: 5px;
                margin: 1px;
            }}
        """)
        layout.addWidget(self.overall_progress)
        
        # Current File Progress Section
        current_label = QLabel("Current File:")
        current_label.setStyleSheet(f"color: {DarkTheme.TEXT_PRIMARY}; font-weight: 600; font-size: 12px;")
        layout.addWidget(current_label)
        
        self.current_file_progress = QProgressBar()
        self.current_file_progress.setStyleSheet(f"""
            QProgressBar {{
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                text-align: center;
                background-color: {DarkTheme.SURFACE_LIGHT};
                height: 20px;
                font-weight: 500;
                font-size: 11px;
            }}
            QProgressBar::chunk {{
                background-color: {DarkTheme.SUCCESS};
                border-radius: 5px;
                margin: 1px;
            }}
        """)
        layout.addWidget(self.current_file_progress)
        
        # Status text
        self.status_label = QLabel("Ready to process")
        self.status_label.setStyleSheet(f"""
            color: {DarkTheme.TEXT_SECONDARY}; 
            font-size: 11px; 
            font-style: italic;
            padding: 3px 0;
        """)
        layout.addWidget(self.status_label)
        
        # Files list
        self.files_list = QListWidget()
        self.files_list.setMaximumHeight(120)  # Reduced height
        self.files_list.setStyleSheet(f"""
            QListWidget {{
                background-color: {DarkTheme.SURFACE};
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                color: {DarkTheme.TEXT_PRIMARY};
                padding: 6px;
                font-size: 11px;
            }}
            QListWidget::item {{
                padding: 4px 8px;
                border-bottom: 1px solid {DarkTheme.BORDER};
                border-radius: 3px;
                margin: 1px 0;
            }}
            QListWidget::item:selected {{
                background-color: {DarkTheme.ACCENT};
                color: white;
            }}
            QListWidget::item:hover:!selected {{
                background-color: {DarkTheme.SURFACE_LIGHT};
            }}
        """)
        layout.addWidget(self.files_list)
        
        self.setLayout(layout)

class FixedMainWindow(QMainWindow):
    """Fixed main window with working functionality"""
    
    def __init__(self):
        super().__init__()
        self.current_files = []
        #self.brand_keywords = []
        self.output_directory = ""
        self.detection_engine = None
        self.document_modifier = None
        self.report_generator = None
        self.performance_processor = None
        self.processing_results = []
        self.setup_ui()
        self.initialize_engine()
        self.apply_dark_theme()
    
    def setup_ui(self):
        self.setWindowTitle("Cloak & Style - PII Data Scrubber")
        self.setGeometry(100, 100, 1400, 900)
        
        # Create central widget
        central_widget = QWidget()
        self.setCentralWidget(central_widget)
        
        # Main layout with proper spacing
        main_layout = QHBoxLayout()
        main_layout.setSpacing(15)  # Reduced spacing
        main_layout.setContentsMargins(15, 15, 15, 15)  # Reduced margins
        
        # Left panel (Input and Options)
        left_panel = QWidget()
        left_layout = QVBoxLayout()
        left_layout.setSpacing(15)  # Reduced spacing
        
        # Header with proper spacing
        header_layout = QHBoxLayout()
        header_layout.setSpacing(10)
        
        title_label = QLabel("Cloak & Style PII Data Scrubber Tool")
        title_label.setFont(QFont("Segoe UI", 14, QFont.Bold))  # Reduced size
        title_label.setStyleSheet(f"color: {DarkTheme.TEXT_PRIMARY}; margin-bottom: 3px;")
        header_layout.addWidget(title_label)
        
        header_layout.addStretch()
        
        # Remove the security badge that was causing visual clutter
        # security_badge = QLabel("🔒 No data sent to cloud • All processing local")
        # security_badge.setStyleSheet(f"""
        #     color: {DarkTheme.WARNING};
        #     font-size: 10px;
        #     font-weight: 500;
        #     padding: 4px 8px;
        #     border: 1px solid {DarkTheme.WARNING};
        #     border-radius: 4px;
        #     background-color: rgba(255, 140, 0, 0.1);
        # """)
        # header_layout.addWidget(security_badge)
        
        left_layout.addLayout(header_layout)
        
        # Input Files section
        input_group = QGroupBox("Input Files")
        input_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 13px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 8px;
                margin-top: 15px;
                padding-top: 15px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 15px;
                padding: 0 8px 0 8px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        input_layout = QVBoxLayout()
        input_layout.setSpacing(15)  # Simple spacing
        
        # Single upload box - clean and simple
        upload_box = QWidget()
        upload_box.setMinimumHeight(100)
        upload_box.setStyleSheet(f"""
            QWidget {{
                background-color: {DarkTheme.SURFACE};
                border: 2px dashed {DarkTheme.BORDER};
                border-radius: 8px;
                padding: 20px;
            }}
        """)
        
        upload_layout = QVBoxLayout()
        upload_layout.setContentsMargins(0, 0, 0, 0)
        upload_layout.setSpacing(10)
        
        # Simple upload label
        upload_label = QLabel("Drag and drop your files here")
        upload_label.setAlignment(Qt.AlignCenter)
        upload_label.setFont(QFont("Segoe UI", 12))
        upload_label.setStyleSheet(f"color: {DarkTheme.TEXT_PRIMARY}; font-weight: 500;")
        upload_layout.addWidget(upload_label)
        
        upload_box.setLayout(upload_layout)
        upload_box.setAcceptDrops(True)
        
        # Connect drag and drop events
        upload_box.dragEnterEvent = lambda event: self.handle_drag_enter(event, upload_box)
        upload_box.dragLeaveEvent = lambda event: self.handle_drag_leave(event, upload_box)
        upload_box.dropEvent = lambda event: self.handle_drop(event, upload_box)
        
        input_layout.addWidget(upload_box)
        
        # Browse button
        browse_btn = QPushButton("Browse Files")
        browse_btn.clicked.connect(self.browse_files)
        browse_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {DarkTheme.ACCENT};
                color: white;
                border: none;
                border-radius: 6px;
                padding: 10px 16px;
                font-weight: 600;
                font-size: 12px;
                min-height: 22px;
            }}
            QPushButton:hover {{
                background-color: #106ebe;
            }}
        """)
        input_layout.addWidget(browse_btn)
        
        # File type note
        file_types_note = QLabel("Supported file types: DOCX, PPTX, PDF, XLSX, CSV, TXT, MD")
        file_types_note.setStyleSheet(f"""
            color: {DarkTheme.TEXT_SECONDARY};
            font-size: 11px;
            font-style: italic;
            padding: 5px 0;
            text-align: center;
        """)
        file_types_note.setAlignment(Qt.AlignCenter)
        input_layout.addWidget(file_types_note)
        
        input_group.setLayout(input_layout)
        left_layout.addWidget(input_group)
        
        # Output Directory section
        output_group = QGroupBox("Output Directory")
        output_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 13px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 8px;
                margin-top: 15px;
                padding-top: 15px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 15px;
                padding: 0 8px 0 8px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        output_layout = QHBoxLayout()
        output_layout.setSpacing(8)
        
        self.output_path = QLineEdit()
        self.output_path.setPlaceholderText("Select output directory (required)")
        self.output_path.setStyleSheet(f"""
            QLineEdit {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                padding: 8px 12px;
                color: {DarkTheme.TEXT_PRIMARY};
                font-size: 11px;
                min-height: 22px;
            }}
            QLineEdit:focus {{
                border-color: {DarkTheme.ACCENT};
            }}
        """)
        output_layout.addWidget(self.output_path)
        
        output_browse_btn = QPushButton("Browse")
        output_browse_btn.clicked.connect(self.browse_output_directory)
        output_browse_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                color: {DarkTheme.TEXT_PRIMARY};
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 6px;
                padding: 8px 16px;
                font-weight: 600;
                font-size: 11px;
                min-height: 22px;
            }}
            QPushButton:hover {{
                background-color: {DarkTheme.BORDER};
            }}
        """)
        output_layout.addWidget(output_browse_btn)
        
        output_group.setLayout(output_layout)
        left_layout.addWidget(output_group)
        
        # Fixed Processing Options
        self.options_widget = FixedProcessingOptions()
        left_layout.addWidget(self.options_widget)
        
        left_panel.setLayout(left_layout)
        left_panel.setMaximumWidth(400)  # Reduced width
        main_layout.addWidget(left_panel)
        
        # Right panel (Preview and Processing)
        right_panel = QWidget()
        right_layout = QVBoxLayout()
        right_layout.setSpacing(15)
        
        # File Preview
        preview_group = QGroupBox("File Preview")
        preview_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 13px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 8px;
                margin-top: 15px;
                padding-top: 15px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 15px;
                padding: 0 8px 0 8px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        preview_layout = QVBoxLayout()
        
        self.preview_widget = FixedFilePreview()
        preview_layout.addWidget(self.preview_widget)
        
        preview_group.setLayout(preview_layout)
        right_layout.addWidget(preview_group)
        
        # Processing section
        processing_group = QGroupBox("Processing")
        processing_group.setStyleSheet(f"""
            QGroupBox {{
                color: {DarkTheme.TEXT_PRIMARY};
                font-weight: 600;
                font-size: 13px;
                border: 1px solid {DarkTheme.BORDER};
                border-radius: 8px;
                margin-top: 15px;
                padding-top: 15px;
                background-color: {DarkTheme.SURFACE};
            }}
            QGroupBox::title {{
                subcontrol-origin: margin;
                left: 15px;
                padding: 0 8px 0 8px;
                background-color: {DarkTheme.SURFACE};
            }}
        """)
        processing_layout = QVBoxLayout()
        processing_layout.setSpacing(10)
        
        self.status_widget = FixedProcessingStatus()
        processing_layout.addWidget(self.status_widget)
        
        # Action buttons
        button_layout = QHBoxLayout()
        button_layout.addStretch()
        
        self.start_btn = QPushButton("Start Processing")
        self.start_btn.clicked.connect(self.start_processing)
        self.start_btn.setEnabled(False)
        self.start_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {DarkTheme.SUCCESS};
                color: white;
                border: none;
                border-radius: 6px;
                padding: 12px 20px;
                font-weight: 600;
                font-size: 12px;
                min-height: 25px;
            }}
            QPushButton:hover {{
                background-color: #0e6e0e;
            }}
            QPushButton:disabled {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                color: {DarkTheme.TEXT_SECONDARY};
            }}
        """)
        button_layout.addWidget(self.start_btn)
        
        self.cancel_btn = QPushButton("Cancel")
        self.cancel_btn.clicked.connect(self.cancel_processing)
        self.cancel_btn.setEnabled(False)
        self.cancel_btn.setStyleSheet(f"""
            QPushButton {{
                background-color: {DarkTheme.ERROR};
                color: white;
                border: none;
                border-radius: 6px;
                padding: 12px 20px;
                font-weight: 600;
                font-size: 12px;
                min-height: 25px;
            }}
            QPushButton:hover {{
                background-color: #b32d30;
            }}
            QPushButton:disabled {{
                background-color: {DarkTheme.SURFACE_LIGHT};
                color: {DarkTheme.TEXT_SECONDARY};
            }}
        """)
        button_layout.addWidget(self.cancel_btn)
        
        processing_layout.addLayout(button_layout)
        processing_group.setLayout(processing_layout)
        right_layout.addWidget(processing_group)
        
        right_panel.setLayout(right_layout)
        main_layout.addWidget(right_panel)
        
        central_widget.setLayout(main_layout)
        
        # Enable drag and drop
        self.setAcceptDrops(True)
    
    def apply_dark_theme(self):
        """Apply dark theme to the entire application"""
        self.setStyleSheet(f"""
            QMainWindow {{
                background-color: {DarkTheme.BACKGROUND};
                color: {DarkTheme.TEXT_PRIMARY};
            }}
            QWidget {{
                background-color: {DarkTheme.BACKGROUND};
                color: {DarkTheme.TEXT_PRIMARY};
            }}
        """)
    
    def initialize_engine(self):
        """Initialize the detection engine"""
        try:
            self.detection_engine = PIIDetectionEngine()
            self.document_modifier = DocumentModifier(self.detection_engine)
            self.report_generator = ReportGenerator()
            self.performance_processor = LaptopOptimizedProcessor()
            print("✅ Detection engine initialized successfully")
        except Exception as e:
            print(f"⚠️ Error initializing detection engine: {e}")
            QMessageBox.warning(self, "Warning", 
                              f"Detection engine initialization failed: {e}\n"
                              "Some features may not work properly.")
    

    def browse_files(self):
        """Browse for files"""
        if self.options_widget.process_subfolders.isChecked():
            files = QFileDialog.getExistingDirectory(
                self,
                "Select Files or Folders",
                ""
            )
            file_list = []
            for root, _ , files in os.walk(files):
                for file in files:
                    if file.endswith(('.txt', '.csv', '.docx', '.pptx', '.xlsx', '.pdf', '.md', '.log')):
                        file_list.append(os.path.join(root, file))
            
            self.current_files = []
            self.add_files(file_list)
        else:
            files, _ = QFileDialog.getOpenFileNames(
                self,
                "Select Files",
                "",
                "All Files (*.txt *.csv *.docx *.pptx *.xlsx *.pdf *.md *.log)"
            )
            if files:
                self.add_files(files)
    
    def browse_output_directory(self):
        """Browse for output directory"""
        directory = QFileDialog.getExistingDirectory(
            self,
            "Select Output Directory",
            ""
        )
        
        if directory:
            self.output_directory = directory
            self.output_path.setText(directory)
            self.update_start_button()
    
    def add_files(self, files):
        """Add files to the list and load preview"""
        # Clear existing files first to avoid duplicates
        if not self.options_widget.process_subfolders.isChecked():
            self.current_files = []
        self.current_files.extend(files)
        self.update_file_list()
        self.update_start_button()
        
        # Load preview of the first file
        if self.current_files:
            self.load_file_preview(self.current_files[0])
        
        print(f"Added {len(files)} files. Total files: {len(self.current_files)}")
        print(f"Output directory: {self.output_directory}")
        print(f"Start button should be enabled: {len(self.current_files) > 0 and self.output_directory != ''}")
    
    def load_file_preview(self, file_path):
        """Load file content into preview"""
        if hasattr(self.preview_widget, 'load_file_content'):
            self.preview_widget.load_file_content(file_path)
    
    def update_file_list(self):
        """Update the file list display"""
        self.status_widget.files_list.clear()
        for file_path in self.current_files:
            item = QListWidgetItem(os.path.basename(file_path))
            item.setToolTip(file_path)
            self.status_widget.files_list.addItem(item)
    
    def update_start_button(self):
        """Update start button state"""
        has_files = len(self.current_files) > 0
        has_output = self.output_directory != ""
        should_enable = has_files and has_output
        
        print(f"update_start_button: files={has_files}, output={has_output}, enable={should_enable}")
        
        self.start_btn.setEnabled(should_enable)
        
        # Update button styling based on state
        if should_enable:
            self.start_btn.setStyleSheet(f"""
                QPushButton {{
                    background-color: {DarkTheme.SUCCESS};
                    color: white;
                    border: none;
                    border-radius: 6px;
                    padding: 12px 20px;
                    font-weight: 600;
                    font-size: 12px;
                    min-height: 25px;
                }}
                QPushButton:hover {{
                    background-color: #0e6e0e;
                }}
            """)
        else:
            self.start_btn.setStyleSheet(f"""
                QPushButton {{
                    background-color: {DarkTheme.SURFACE_LIGHT};
                    color: {DarkTheme.TEXT_SECONDARY};
                    border: none;
                    border-radius: 6px;
                    padding: 12px 20px;
                    font-weight: 600;
                    font-size: 12px;
                    min-height: 25px;
                }}
            """)
    
    def handle_drag_enter(self, event, widget):
        """Handle drag enter event"""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
            widget.setStyleSheet(f"""
                QWidget {{
                    background-color: {DarkTheme.HIGHLIGHT};
                    border: 2px dashed {DarkTheme.ACCENT};
                    border-radius: 8px;
                    padding: 20px;
                }}
            """)
    
    def handle_drag_leave(self, event, widget):
        """Handle drag leave event"""
        widget.setStyleSheet(f"""
            QWidget {{
                background-color: {DarkTheme.SURFACE};
                border: 2px dashed {DarkTheme.BORDER};
                border-radius: 8px;
                padding: 20px;
            }}
        """)
    
    def handle_drop(self, event, widget):
        """Handle drop event"""
        files = []
        for url in event.mimeData().urls():
            file_path = url.toLocalFile()
            if os.path.isfile(file_path):
                files.append(file_path)
        
        if files:
            self.add_files(files)
        
        # Reset styling
        widget.setStyleSheet(f"""
            QWidget {{
                background-color: {DarkTheme.SURFACE};
                border: 2px dashed {DarkTheme.BORDER};
                border-radius: 8px;
                padding: 20px;
            }}
        """)
    
    def start_processing(self):
        """Start processing files"""
        if not self.current_files:
            QMessageBox.warning(self, "No Files", "Please select files to process.")
            return
        
        if not self.output_directory:
            QMessageBox.warning(self, "No Output Directory", "Please select an output directory.")
            return
        
        # Start processing logic here
        self.start_btn.setEnabled(False)
        self.cancel_btn.setEnabled(True)
        self.status_widget.status_label.setText("Processing...")
        self.status_widget.overall_progress.setValue(0)
        self.status_widget.current_file_progress.setValue(0)
        
        try:
            # Process each file
            total_files = len(self.current_files)
            for i, file_path in enumerate(self.current_files):
                # Update progress
                overall_progress = int((i / total_files) * 100)
                self.status_widget.overall_progress.setValue(overall_progress)
                self.status_widget.status_label.setText(f"Processing: {os.path.basename(file_path)}")
                # Process current file
                self.process_single_file(file_path)
                # Update current file progress
                self.status_widget.current_file_progress.setValue(100)
                # Small delay to show progress
                QApplication.processEvents()
            # Complete
            self.status_widget.overall_progress.setValue(100)
            self.status_widget.status_label.setText("Processing completed successfully")
            
            if not self.options_widget.dry_run.isChecked():
                QMessageBox.information(self, "Processing Complete", 
                                      f"Successfully processed {len(self.current_files)} file(s).\n"
                                      f"Output saved to: {self.output_directory}")
            else:
                QMessageBox.information(self, "Processing Complete", 
                                      f"Successfully processed {len(self.current_files)} file(s).\n"
                                      "Go to Detected Entities to review the results.")
            
        except Exception as e:
            QMessageBox.critical(self, "Processing Error", 
                               f"An error occurred during processing:\n{str(e)}")
            self.status_widget.status_label.setText("Processing failed")
        
        finally:
            self.start_btn.setEnabled(True)
            self.cancel_btn.setEnabled(False)
    
    def process_single_file(self, file_path):
        """Process a single file for PII detection and masking"""
        try:
            # Read file content
            if file_path.endswith(('.csv', '.xlsx')):
                # For CSV and Excel, read as table
                import pandas as pd
                df = pd.read_csv(file_path) if file_path.endswith('.csv') else pd.read_excel(file_path)
                content = df.to_string()
            elif file_path.endswith(('.docx', '.pptx', '.pdf')):
                from docx import Document
                from pptx import Presentation
                import PyPDF2
                if file_path.endswith('.docx'):
                    doc = Document(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs])
                elif file_path.endswith('.pptx'):
                    prs = Presentation(file_path)
                    content = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content += shape.text + "\n"
                elif file_path.endswith('.pdf'):
                    content = ""
                    with open(file_path, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        for page in reader.pages:
                            content += page.extract_text() + "\n"
            else:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()

            # Detect PII entities (if detection engine is available)
            detected_entities = []
            if self.detection_engine and hasattr(self.detection_engine, 'detect_pii'):
                try:
                    detected_entities = self.detection_engine.detect_pii(content)
                except Exception as e:
                    print(f"Detection engine error: {e}")
            
            # Create masked content (simple masking for demo)
            masked_content = self.mask_pii_content(content, detected_entities, self.options_widget.mask_format.currentText(), brand_keywords)
            # Save processed file
            output_filename = f"masked_{os.path.basename(file_path)}"
            output_path = os.path.join(self.output_directory, output_filename)
            
            if not self.options_widget.dry_run.isChecked():
                if file_path.endswith(('.csv', '.xlsx')):
                    # Preserve structure for XLSX; for CSV, write text directly
                    if file_path.endswith('.csv'):
                        with open(output_path, 'w', encoding='utf-8', newline='') as f:
                            f.write(masked_content)
                    elif file_path.endswith('.xlsx'):
                        try:
                            from openpyxl import load_workbook
                            wb = load_workbook(filename=file_path)
                            for ws in wb.worksheets:
                                for row in ws.iter_rows():
                                    for cell in row:
                                        try:
                                            value = cell.value
                                            if isinstance(value, str) and value:
                                                new_value = self.mask_pii_content(value, detected_entities, self.options_widget.mask_format.currentText(), brand_keywords)
                                                if new_value != value:
                                                    cell.value = new_value
                                        except Exception:
                                            continue
                            wb.save(output_path)
                        except Exception:
                            # Fallback to pandas if openpyxl path fails
                            import pandas as pd
                            from io import StringIO
                            try:
                                df_masked = pd.read_csv(StringIO(masked_content))
                                df_masked.to_excel(output_path, index=False)
                            except Exception:
                                # Last resort: write masked text as .txt with .xlsx extension to avoid crashing
                                with open(output_path, 'w', encoding='utf-8') as f:
                                    f.write(masked_content)
                elif file_path.endswith('.docx'):
                    from docx import Document
                    doc = Document(file_path)
                    # Mask paragraphs
                    for para in doc.paragraphs:
                        para.text = self.mask_pii_content(para.text, detected_entities, self.options_widget.mask_format.currentText(), brand_keywords)
                    # Mask tables
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                cell.text = self.mask_pii_content(cell.text, detected_entities, self.options_widget.mask_format.currentText(), brand_keywords)
                    doc.save(output_path)
                elif file_path.endswith('.pptx'):
                    from pptx import Presentation
                    prs = Presentation(file_path)

                    # Helper to mask text in a text_frame while preserving paragraphs/runs formatting
                    def mask_text_frame(text_frame):
                        for paragraph in text_frame.paragraphs:
                            for run in paragraph.runs:
                                try:
                                    if run.text:
                                        new_text = self.mask_pii_content(run.text, detected_entities, self.options_widget.mask_format.currentText(), brand_keywords)
                                        if new_text != run.text:
                                            run.text = new_text
                                except Exception:
                                    # Skip problematic runs, continue with others
                                    continue

                    for slide in prs.slides:
                        for shape in slide.shapes:
                            # Text-containing shapes
                            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                                mask_text_frame(shape.text_frame)
                            # Table cells
                            if hasattr(shape, "has_table") and shape.has_table:
                                table = shape.table
                                for row in table.rows:
                                    for cell in row.cells:
                                        if hasattr(cell, "text_frame") and cell.text_frame:
                                            mask_text_frame(cell.text_frame)

                    prs.save(output_path)
                elif file_path.endswith('.pdf'):
                    # Replace PII text with masked text via redaction annotations (with replacement text) to preserve layout
                    import fitz
                    try:
                        doc = fitz.open(file_path)
                        has_entities = bool(self.detection_engine and detected_entities and hasattr(detected_entities, 'entities_found'))
                        if has_entities:
                            for page in doc:
                                for entity in getattr(detected_entities, 'entities_found', []) or []:
                                    value = str(getattr(entity, 'value', '') or '')
                                    if not value:
                                        continue
                                    replacement = self.mask_pii_content(value, detected_entities, self.options_widget.mask_format.currentText(), brand_keywords)
                                    if replacement == value:
                                        continue
                                    for rect in page.search_for(value):
                                        # Add redaction without fill/border; overlay masked text to avoid visible box
                                        try:
                                            page.add_redact_annot(
                                                rect,
                                                fill=None,
                                                text=replacement,
                                                text_color=(0, 0, 0),
                                                align=0,
                                                overlay=True,
                                                border=None
                                            )
                                        except TypeError:
                                            # Fallback for older PyMuPDF without overlay/border args
                                            page.add_redact_annot(rect, fill=None, text=replacement, text_color=(0, 0, 0), align=0)
                                # Apply redactions for the page so the underlying text is actually replaced/removed
                                page.apply_redactions()
                            doc.save(output_path)
                            doc.close()
                        else:
                            # Fallback: rebuild a Unicode-safe PDF with text-only when no entities available
                            doc2 = fitz.open()
                            page2 = doc2.new_page()
                            x, y = 50, 50
                            line_height = 14
                            max_y = page2.rect.height - 50
                            for line in masked_content.split('\n'):
                                if y > max_y:
                                    page2 = doc2.new_page()
                                    y = 50
                                page2.insert_text((x, y), line, fontsize=12)
                                y += line_height
                            doc2.save(output_path)
                            doc2.close()
                    except Exception as e:
                        # Fallback: rebuild a Unicode-safe PDF if redaction replacement fails
                        print(e)
                        doc = fitz.open()
                        page = doc.new_page()
                        x, y = 50, 50
                        line_height = 14
                        max_y = page.rect.height - 50
                        for line in masked_content.split('\n'):
                            if y > max_y:
                                page = doc.new_page()
                                y = 50
                            page.insert_text((x, y), line, fontsize=12)
                            y += line_height
                        doc.save(output_path)
                        doc.close()
                else:
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(masked_content)
                # Update preview with detected entities

            if detected_entities:
                entities_text = "Detected PII Entities:\n\n"
                for entity in detected_entities.entities_found:
                    entities_text += f"• {entity.entity_type}: {entity.value}\n"
                self.preview_widget.entities_preview.setText(entities_text)
            print(f"Processed: {file_path} -> {output_path}")
            
        except Exception as e:
            print(f"Error processing {file_path}: {e}")
            raise
    
    def mask_pii_content(self, content, entities, mask_format, keywords=None):
        """Mask PII content based on detected entities"""
        masked_content = content
        
        # Simple masking patterns (in a real implementation, this would be more sophisticated)
        import re

        # Email addresses
        match (mask_format):
            case 'Token Format [TYPE_###]':
                masked_content = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[EMAIL]', masked_content)
            case 'Asterisk Format [***]':
                masked_content = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[******]', masked_content)
        #masked_content = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '[EMAIL]', masked_content)
        
        # Phone numbers (various formats)
        match (mask_format):
            case 'Token Format [TYPE_###]':
                masked_content = re.sub(r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', '[PHONE]', masked_content)
            case 'Asterisk Format [***]':
                masked_content = re.sub(r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b', '[******]', masked_content)
        
        match (mask_format):
            case 'Token Format [TYPE_###]':
                masked_content = re.sub(r'\b\(\d{3}\)\s*\d{3}[-.]?\d{4}\b', '[PHONE]', masked_content)
            case 'Asterisk Format [***]':
                masked_content = re.sub(r'\b\(\d{3}\)\s*\d{3}[-.]?\d{4}\b', '[******]', masked_content)
        
        # Social Security Numbers
        match (mask_format):
            case 'Token Format [TYPE_###]':
                masked_content = re.sub(r'\b\d{3}-\d{2}-\d{4}\b', '[SSN]', masked_content)
            case 'Asterisk Format [***]':
                masked_content = re.sub(r'\b\d{3}-\d{2}-\d{4}\b', '[******]', masked_content)
        
        # Credit Card Numbers
        match (mask_format):
            case 'Token Format [TYPE_###]':
                masked_content = re.sub(r'\b\d{4}[- ]?\d{4}[- ]?\d{4}[- ]?\d{4}\b', '[CREDIT_CARD]', masked_content)
            case 'Asterisk Format [***]':
                masked_content = re.sub(r'\b\d{4}[- ]?\d{4}[- ]?\d{4}[- ]?\d{4}\b', '[******]', masked_content)
        
        # IP Addresses
        match (mask_format):
            case 'Token Format [TYPE_###]':
                masked_content = re.sub(r'\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b', '[IP_ADDRESS]', masked_content)
            case 'Asterisk Format [***]':
                masked_content = re.sub(r'\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b', '[******]', masked_content)
        
        if keywords:
            for brand in keywords:
                masked_content = re.sub(brand, '[BRAND]', masked_content)
            

        return masked_content
    
    def cancel_processing(self):
        """Cancel processing"""
        self.start_btn.setEnabled(True)
        self.cancel_btn.setEnabled(False)
        self.status_widget.status_label.setText("Processing cancelled")
        self.status_widget.overall_progress.setValue(0)
        self.status_widget.current_file_progress.setValue(0)

def main():
    app = QApplication(sys.argv)
    
    # Set application style
    app.setStyle('Fusion')
    
    # Create and show main window
    window = FixedMainWindow()
    window.show()
    
    sys.exit(app.exec())

if __name__ == "__main__":
    main()
